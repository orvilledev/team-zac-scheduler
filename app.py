from flask import Flask, render_template, request, redirect, url_for, flash, send_from_directory, abort, jsonify, make_response, Response  # type: ignore
from flask_login import LoginManager, login_user, logout_user, login_required, current_user  # type: ignore
from flask_wtf.csrf import generate_csrf, CSRFProtect  # type: ignore
from werkzeug.utils import secure_filename  # type: ignore
from functools import wraps
import os
import tempfile
from datetime import datetime, date, timedelta
from apscheduler.schedulers.background import BackgroundScheduler  # type: ignore
from apscheduler.triggers.date import DateTrigger  # type: ignore
import atexit

from config import Config
from models import db, User, Musician, SundayService, ServiceMusician, Practice, PracticeMusician, Song, MusicianAvailability, Slide, ProfilePost, PracticeSong, PostLike, PostHeart, PostRepost, PostComment, EventAnnouncement, Notification, SMSLog, UserPermission, Journal, LeaveRequest, ActivityLog, Task, TaskOption, Tool, Message
from sms_service import send_practice_assignment_sms, send_practice_reminder_sms, format_phone_number
from forms import (
    LoginForm, MusicianForm, ServiceForm, ServiceMusicianForm,
    PracticeForm, PracticeMusicianForm, UserForm, SlideForm, ProfilePostForm, PostCommentForm, EventAnnouncementForm, ProfileCustomizationForm, PermissionForm, JournalForm, ToolForm
)

app = Flask(__name__)
app.config.from_object(Config)
csrf = CSRFProtect(app)

# Initialize APScheduler for SMS reminders
# Only start scheduler once, not on every reload in debug mode
scheduler = None
if not app.debug or os.environ.get('WERKZEUG_RUN_MAIN') == 'true':
    try:
        scheduler = BackgroundScheduler()
        if not scheduler.running:
            scheduler.start()
        atexit.register(lambda: scheduler.shutdown() if scheduler and scheduler.running else None)
    except Exception as e:
        print(f"Warning: Could not start scheduler: {e}")
        scheduler = None

# Make CSRF token helper available in all templates
@app.context_processor
def inject_csrf_token():
    def get_csrf_token():
        return generate_csrf()
    
    # Also inject unread notification count for authenticated users
    unread_notification_count = 0
    if current_user.is_authenticated:
        unread_notification_count = Notification.query.filter_by(user_id=current_user.id, is_read=False).count()
    
    return dict(
        get_csrf_token=get_csrf_token,
        unread_notification_count=unread_notification_count
    )

# Jinja2 filter to format instrument names
@app.template_filter('format_instrument')
def format_instrument(instrument):
    """Format instrument name with appropriate suffix"""
    if not instrument:
        return ""
    
    instrument_lower = instrument.lower().strip()
    
    # Special cases
    if instrument_lower == 'drums' or instrument_lower == 'drum':
        return 'Drummer'
    elif instrument_lower == 'vocals' or instrument_lower == 'vocal':
        return 'Vocalist'
    elif instrument_lower == 'keyboard' or instrument_lower == 'keyboards':
        return 'Keyboardist'
    else:
        # Add "player" suffix for other instruments
        return f"{instrument} player"

# Jinja2 filter to convert UTC datetime to Manila time
@app.template_filter('bold_title')
def bold_title_filter(text, title):
    """Bold the title if it appears at the start of the text"""
    if not text or not title:
        return text
    # Remove all leading whitespace
    text = text.lstrip()
    title = title.strip()
    if text.lower().startswith(title.lower()):
        # Preserve the original spacing after the title
        remaining_text = text[len(title):]
        return f'<strong>{text[:len(title)]}</strong>{remaining_text}'
    return text

@app.template_filter('manila_time')
def manila_time(dt):
    """Convert UTC datetime to Manila time (UTC+8)"""
    if not dt:
        return None
    
    try:
        # Try using pytz if available
        import pytz  # type: ignore
        utc = pytz.UTC
        manila_tz = pytz.timezone('Asia/Manila')
        
        # If datetime is naive, assume it's UTC
        if dt.tzinfo is None:
            dt = utc.localize(dt)
        
        # Convert to Manila time
        manila_dt = dt.astimezone(manila_tz)
        return manila_dt
    except ImportError:
        # Fallback: manually add 8 hours if pytz is not available
        from datetime import timedelta, timezone
        if dt.tzinfo is None:
            # Assume UTC if naive, add 8 hours
            return dt + timedelta(hours=8)
        else:
            # If it has timezone info, convert to UTC first, then add 8 hours
            # Convert to UTC
            utc_dt = dt.astimezone(timezone.utc) if hasattr(dt, 'astimezone') else dt
            # Add 8 hours for Manila time
            return utc_dt.replace(tzinfo=None) + timedelta(hours=8)

# Helper function to format Manila time (available in templates)
@app.template_filter('format_manila_time')
def format_manila_time_filter(dt, format_string=None):
    """Format datetime in Manila time - Jinja2 filter
    
    Usage: {{ datetime|format_manila_time('%B %d, %Y at %I:%M %p') }}
    """
    if not dt:
        return ''
    if format_string is None:
        format_string = '%B %d, %Y at %I:%M %p'
    manila_dt = manila_time(dt)
    if manila_dt:
        return manila_dt.strftime(format_string)
    return ''

@app.context_processor
def inject_manila_time_formatter():
    """Make manila_time formatter available in templates as function (for function calls)"""
    def format_manila_time_func(dt, format_string='%B %d, %Y at %I:%M %p'):
        """Format datetime in Manila time"""
        if not dt:
            return ''
        manila_dt = manila_time(dt)
        if manila_dt:
            return manila_dt.strftime(format_string)
        return ''
    return dict(format_manila_time_func=format_manila_time_func)

# Initialize extensions
db.init_app(app)
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = 'Please log in to access this page.'
login_manager.login_message_category = 'info'


@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))


# Role-based decorators
def admin_required(f):
    @wraps(f)
    @login_required
    def decorated_function(*args, **kwargs):
        if not current_user.is_admin():
            flash('Admin or Team Leader access required.', 'danger')
            return redirect(url_for('dashboard'))
        return f(*args, **kwargs)
    return decorated_function


def worship_leader_required(f):
    @wraps(f)
    @login_required
    def decorated_function(*args, **kwargs):
        if not current_user.is_worship_leader():
            flash('Worship leader or admin access required.', 'danger')
            return redirect(url_for('dashboard'))
        return f(*args, **kwargs)
    return decorated_function


def permission_required(permission_type):
    """Decorator to check if user has a specific permission"""
    def decorator(f):
        @wraps(f)
        @login_required
        def decorated_function(*args, **kwargs):
            if not (current_user.is_admin() or current_user.is_worship_leader() or current_user.has_permission(permission_type)):
                flash('You do not have permission to perform this action.', 'danger')
                return redirect(url_for('dashboard'))
            return f(*args, **kwargs)
        return decorated_function
    return decorator


# Authentication routes
@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    
    form = LoginForm()
    if form.validate_on_submit():
        username = form.username.data.strip() if form.username.data else ''
        
        if not username:
            flash('Please enter a username.', 'danger')
            return render_template('login.html', form=form)
        
        # Find user by username (case-sensitive)
        user = User.query.filter_by(username=username).first()
        
        if not user:
            flash('User not found.', 'danger')
            return render_template('login.html', form=form)
        
        # Login without password check (passwords disabled)
        login_user(user, remember=True)
        next_page = request.args.get('next')
        return redirect(next_page) if next_page else redirect(url_for('dashboard'))
    
    return render_template('login.html', form=form)


@app.route('/logout')
@login_required
def logout():
    try:
        logout_user()
        flash('You have been logged out.', 'info')
        return redirect(url_for('login'))
    except Exception as e:
        # If logout fails, still try to redirect
        print(f"Error during logout: {e}")
        return redirect(url_for('login'))


# Helper function to log activities
def log_activity(activity_type, actor_id, description, target_user_id=None, slide_id=None, leave_request_id=None, metadata=None):
    """Log an activity to the activity log"""
    try:
        import json
        activity = ActivityLog(
            activity_type=activity_type,
            actor_id=actor_id,
            target_user_id=target_user_id,
            description=description,
            slide_id=slide_id,
            leave_request_id=leave_request_id,
            extra_data=json.dumps(metadata) if metadata else None
        )
        db.session.add(activity)
        db.session.commit()
    except Exception as e:
        # Don't fail the main operation if logging fails
        print(f"Error logging activity: {e}")
        db.session.rollback()


# Dashboard
@app.route('/')
@app.route('/dashboard')
@login_required
def dashboard():
    # Get upcoming services (next 4 weeks)
    today = date.today()
    upcoming_services = SundayService.query.filter(
        SundayService.date >= today
    ).order_by(SundayService.date).limit(5).all()
    
    # Get only the latest upcoming practice with eager loading
    from sqlalchemy.orm import joinedload  # type: ignore
    from models import PracticeMusician, PracticeSong
    latest_practice = Practice.query.options(
        joinedload(Practice.musicians).joinedload(PracticeMusician.musician),
        joinedload(Practice.songs).joinedload(PracticeSong.song),
        joinedload(Practice.songs).joinedload(PracticeSong.preparer)
    ).filter(
        Practice.date >= today
    ).order_by(Practice.date).first()
    
    # Check if current user is assigned to the latest practice
    user_assignment_info = None
    if latest_practice:
        # Get or create team member profile for current user if needed
        if not current_user.musician:
            musician = Musician(
                name=current_user.get_display_name(),
                user_id=current_user.id,
                instruments=current_user.role if current_user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader'] else None
            )
            db.session.add(musician)
            db.session.commit()
            
            # Log activity for new member
            log_activity(
                activity_type='new_member',
                actor_id=current_user.id,
                target_user_id=current_user.id,
                description=f"{current_user.get_display_name()} joined the team as a new member",
                metadata={'role': current_user.role or 'member'}
            )
        elif not current_user.musician.instruments and current_user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
            # Sync role if musician profile exists but doesn't have a role set
            current_user.musician.instruments = current_user.role
            db.session.commit()
            musician = current_user.musician
        else:
            musician = current_user.musician
        
        # Check if user is assigned to this practice
        # Query PracticeMusician directly to ensure we get all assignments
        practice_assignments = PracticeMusician.query.filter_by(practice_id=latest_practice.id).all()
        for assignment in practice_assignments:
            if assignment.musician and assignment.musician_id == musician.id:
                user_assignment_info = {
                    'instrument': assignment.instrument,
                    'nickname': current_user.get_display_name(),
                    'date': latest_practice.date.strftime('%B %d, %Y')
                }
                break
    
    # Get newly added musicians (created within the last 30 days)
    # Only include musicians that have a valid user_id pointing to an existing User
    from datetime import timedelta
    from collections import defaultdict
    from sqlalchemy.orm import joinedload  # type: ignore
    thirty_days_ago = datetime.utcnow() - timedelta(days=30)
    
    # Get all valid user IDs
    valid_user_ids = [user.id for user in User.query.all()]
    
    # Filter to only include musicians with valid user_id that exists in User table
    all_new_musicians = Musician.query.options(
        joinedload(Musician.user)
    ).filter(
        Musician.created_at >= thirty_days_ago,
        Musician.user_id.isnot(None)  # Must have a user_id
    ).order_by(Musician.created_at.desc()).all()
    
    # Filter to ensure user actually exists and is valid
    all_new_musicians = [m for m in all_new_musicians if m.user_id and m.user_id in valid_user_ids and m.user is not None]
    
    # Group musicians by user_id or display name to avoid duplicates and combine instruments
    musician_dict = {}
    
    for musician in all_new_musicians:
        display_name = musician.get_display_name().strip().lower() if musician.get_display_name() else musician.name.strip().lower()
        
        # Find existing entry by user_id or display name
        existing_musician = None
        existing_key = None
        
        # First, check if we already have this user_id
        if musician.user_id:
            for key, existing in musician_dict.items():
                if existing.user_id == musician.user_id:
                    existing_musician = existing
                    existing_key = key
                    break
        
        # If not found by user_id, check by display name
        if not existing_musician:
            for key, existing in musician_dict.items():
                existing_display = existing.get_display_name().strip().lower() if existing.get_display_name() else existing.name.strip().lower()
                if existing_display == display_name:
                    existing_musician = existing
                    existing_key = key
                    break
        
        if existing_musician:
            # Duplicate found - combine instruments
            existing_instruments = set()
            new_instruments = set()
            
            if existing_musician.instruments:
                existing_instruments = {inst.strip() for inst in existing_musician.instruments.split(',') if inst.strip()}
            if musician.instruments:
                new_instruments = {inst.strip() for inst in musician.instruments.split(',') if inst.strip()}
            
            combined_instruments = ', '.join(sorted(existing_instruments | new_instruments))
            existing_musician.instruments = combined_instruments if combined_instruments else None
            
            # Also merge other fields if the existing one is missing them
            if not existing_musician.profile_picture and musician.profile_picture:
                existing_musician.profile_picture = musician.profile_picture
            if not existing_musician.banner and musician.banner:
                existing_musician.banner = musician.banner
        else:
            # New musician - add to dict
            # Use user_id as key if available, otherwise use display name
            if musician.user_id:
                key = f"user_{musician.user_id}"
            else:
                key = f"name_{display_name}"
            musician_dict[key] = musician
    
    # Convert to list - no limit on number of profiles
    new_musicians = list(musician_dict.values())
    
    # Get active event announcements
    announcements = EventAnnouncement.query.filter_by(is_active=True).order_by(EventAnnouncement.display_order, EventAnnouncement.created_at.desc()).all()
    
    # Get pending leave requests for team leaders/admins
    pending_leave_requests = []
    if current_user.is_team_leader():
        pending_leave_requests = LeaveRequest.query.filter_by(status='pending').order_by(LeaveRequest.date.asc()).all()
    
    # Get recent activities (last 30 days, limit 20)
    from datetime import timedelta
    from sqlalchemy.orm import joinedload  # type: ignore
    thirty_days_ago = datetime.utcnow() - timedelta(days=30)
    recent_activities = ActivityLog.query.options(
        joinedload(ActivityLog.actor).joinedload(User.musician),
        joinedload(ActivityLog.slide)
    ).filter(
        ActivityLog.created_at >= thirty_days_ago
    ).order_by(ActivityLog.created_at.desc()).limit(20).all()
    
    return render_template('dashboard.html',
                         upcoming_services=upcoming_services,
                         latest_practice=latest_practice,
                         user_assignment_info=user_assignment_info,
                         new_musicians=new_musicians,
                         announcements=announcements,
                         pending_leave_requests=pending_leave_requests,
                         recent_activities=recent_activities)


# Musician routes
@app.route('/fba-copy')
@login_required
def fba_copy():
    # Show all users, not just those with team member profiles
    users_list = User.query.order_by(User.username).all()
    
    # Sync all musician names with their user display names
    updated_count = 0
    for user in users_list:
        if user.musician and user.musician.name != user.get_display_name():
            user.musician.name = user.get_display_name()
            updated_count += 1
    if updated_count > 0:
        db.session.commit()
    
    return render_template('fba_copy.html', users=users_list)


@app.route('/musicians/add', methods=['GET', 'POST'])
@worship_leader_required
def add_musician():
    form = MusicianForm()
    if form.validate_on_submit():
        musician = Musician(
            name=form.name.data,
            email=form.email.data or None,
            phone=form.phone.data or None,
            mobile=form.mobile.data or None,
            outlook_email=form.outlook_email.data or None,
            whatsapp=form.whatsapp.data or None,
            instruments=form.instruments.data or None,
            bio=form.bio.data or None,
            interests=form.interests.data or None
        )
        db.session.add(musician)
        db.session.commit()
        flash('Team member added successfully.', 'success')
        return redirect(url_for('musicians'))
    return render_template('musician_form.html', form=form, title='Add Team Member')


@app.route('/musicians/user/<int:user_id>/edit', methods=['GET', 'POST'])
@worship_leader_required
def edit_musician_by_user(user_id):
    """Edit team member profile by user_id - creates profile if it doesn't exist"""
    user = User.query.get_or_404(user_id)
    
    # Get or create musician profile
    musician = user.musician
    if not musician:
        # Create a new musician profile for this user, syncing role
        musician = Musician(
            name=user.get_display_name(),
            user_id=user.id,
            instruments=user.role if user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader'] else None
        )
        db.session.add(musician)
        db.session.flush()  # Flush to get the ID
    else:
        # Sync name with user's display name if they differ
        user_display_name = user.get_display_name()
        if musician.name != user_display_name:
            musician.name = user_display_name
        # Sync role if musician profile exists but doesn't have a role set
        if not musician.instruments and user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
            musician.instruments = user.role
        db.session.flush()
    
    form = MusicianForm(obj=musician)
    # If instruments field has old comma-separated data, clear it so user selects from dropdown
    if musician.instruments and musician.instruments not in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
        form.instruments.data = ''
    if form.validate_on_submit():
        musician.name = form.name.data
        musician.email = form.email.data or None
        musician.phone = form.phone.data or None
        musician.mobile = form.mobile.data or None
        musician.outlook_email = form.outlook_email.data or None
        musician.whatsapp = form.whatsapp.data or None
        musician.instruments = form.instruments.data or None
        musician.bio = form.bio.data or None
        musician.interests = form.interests.data or None
        musician.user_id = user.id  # Ensure user_id is set
        db.session.commit()
        flash('Team member profile updated successfully.', 'success')
        return redirect(url_for('musicians'))
    
    return render_template('musician_form.html', form=form, title='Edit Team Member', musician=musician)


@app.route('/musicians/<int:id>/post', methods=['POST'])
@login_required
def create_profile_post(id):
    """Create a new post on a musician's profile wall"""
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Check permissions - users can only post on their own profile
    if not user or current_user.id != user.id:
        flash('You can only post on your own profile.', 'danger')
        return redirect(url_for('view_musician_profile', id=id))
    
    form = ProfilePostForm()
    if form.validate_on_submit():
        post = ProfilePost(musician_id=musician.id)
        post.content = form.content.data or None
        
        # Handle image upload
        if form.image.data:
            file = form.image.data
            # Check if it's a file object (has filename attribute) and not a string
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"post_img_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['POSTS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['POSTS_FOLDER'], filename)
                file.save(file_path)
                post.image_path = f"profiles/posts/{filename}"
        
        # Handle video upload
        if form.video.data:
            file = form.video.data
            # Check if it's a file object (has filename attribute) and not a string
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"post_vid_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['POSTS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['POSTS_FOLDER'], filename)
                file.save(file_path)
                post.video_path = f"profiles/posts/{filename}"
        
        # Ensure at least content, image, or video is provided
        if not post.content and not post.image_path and not post.video_path:
            flash('Please provide text, an image, or a video for your post.', 'warning')
            return redirect(url_for('view_musician_profile', id=id))
        
        db.session.add(post)
        db.session.commit()
        flash('Post created successfully!', 'success')
    
    return redirect(url_for('view_musician_profile', id=id))


@app.route('/profile')
@login_required
def my_profile():
    """Redirect to current user's profile, creating one if needed"""
    # Get or create musician profile for current user
    musician = current_user.musician
    if not musician:
        # Auto-create musician profile if it doesn't exist, syncing role
        musician = Musician(
            name=current_user.get_display_name(),
            user_id=current_user.id,
            instruments=current_user.role if current_user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader'] else None
        )
        db.session.add(musician)
        db.session.commit()
    else:
        # Sync name with user's display name if they differ
        user_display_name = current_user.get_display_name()
        if musician.name != user_display_name:
            musician.name = user_display_name
        # Sync role if musician profile exists but doesn't have a role set
        if not musician.instruments and current_user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
            musician.instruments = current_user.role
        db.session.commit()
    
    return redirect(url_for('view_musician_profile', id=musician.id))


@app.route('/musicians/<int:id>/profile')
@login_required
def view_musician_profile(id):
    """View team member profile in a modern social media style"""
    # Force a fresh query to avoid cached data
    db.session.expire_all()
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Sync musician name with user's display name if they differ
    if user and musician.name != user.get_display_name():
        musician.name = user.get_display_name()
        db.session.commit()
    
    # Get upcoming services for this musician
    from datetime import date
    today = date.today()
    upcoming_services = []
    
    for assignment in musician.service_assignments:
        if assignment.service and assignment.service.date >= today:
            upcoming_services.append({
                'date': assignment.service.date,
                'theme': assignment.service.theme,
                'instrument': assignment.instrument,
                'role': assignment.role
            })
    
    # Sort by date
    upcoming_services.sort(key=lambda x: x['date'])
    
    # Get posts for this musician (most recent first)
    posts = ProfilePost.query.filter_by(musician_id=musician.id).order_by(ProfilePost.created_at.desc()).all()
    
    # Track profile views (only if not viewing own profile)
    if current_user.is_authenticated:
        if not user or current_user.id != user.id:
            musician.profile_views = (musician.profile_views or 0) + 1
            db.session.commit()
    
    # Check if current user can edit this profile
    can_edit = (user and current_user.id == user.id) or current_user.is_worship_leader()
    
    form = ProfilePostForm()
    
    return render_template('musician_profile.html', 
                         musician=musician, 
                         user=user,
                         upcoming_services=upcoming_services[:5],
                         posts=posts,
                         can_edit=can_edit,
                         post_form=form,
                         current_user_id=current_user.id)


@app.route('/musicians/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit_musician(id):
    """Edit team member profile - users can edit their own, worship leaders can edit any"""
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Check permissions
    can_edit = (user and current_user.id == user.id) or current_user.is_worship_leader()
    if not can_edit:
        flash('You do not have permission to edit this profile.', 'danger')
        return redirect(url_for('view_musician_profile', id=id))
    
    form = MusicianForm(obj=musician)
    # If instruments field has old comma-separated data, clear it so user selects from dropdown
    if musician.instruments and musician.instruments not in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
        form.instruments.data = ''
    if form.validate_on_submit():
        musician.name = form.name.data
        musician.email = form.email.data or None
        musician.phone = form.phone.data or None
        musician.mobile = form.mobile.data or None
        musician.outlook_email = form.outlook_email.data or None
        musician.whatsapp = form.whatsapp.data or None
        musician.instruments = form.instruments.data or None
        musician.bio = form.bio.data or None
        musician.interests = form.interests.data or None
        
        # Handle profile picture upload
        if form.profile_picture.data:
            file = form.profile_picture.data
            # Check if it's a file object (has filename attribute) and not a string
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"profile_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['PROFILE_PICTURES_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['PROFILE_PICTURES_FOLDER'], filename)
                file.save(file_path)
                musician.profile_picture = f"profiles/pictures/{filename}"
        
        # Handle banner upload
        if form.banner.data:
            file = form.banner.data
            # Check if it's a file object (has filename attribute) and not a string
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"banner_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['BANNERS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['BANNERS_FOLDER'], filename)
                file.save(file_path)
                musician.banner = f"profiles/banners/{filename}"
        
        db.session.commit()
        flash('Profile updated successfully.', 'success')
        return redirect(url_for('view_musician_profile', id=id))
    return render_template('musician_form.html', form=form, title='Edit Profile', musician=musician)


@app.route('/musicians/<int:id>/edit-picture', methods=['GET', 'POST'])
@login_required
def edit_musician_picture(id):
    """Quick edit for profile picture and banner only"""
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Check permissions
    can_edit = (user and current_user.id == user.id) or current_user.is_worship_leader()
    if not can_edit:
        flash('You do not have permission to edit this profile.', 'danger')
        return redirect(url_for('view_musician_profile', id=id))
    
    form = MusicianForm(obj=musician)
    # Only show profile picture and banner fields
    if request.method == 'POST':
        if form.validate_on_submit():
            # Handle profile picture upload
            if form.profile_picture.data:
                file = form.profile_picture.data
                if hasattr(file, 'filename') and file.filename:
                    # Delete old profile picture if exists
                    if musician.profile_picture:
                        old_file_path = os.path.join(app.static_folder, musician.profile_picture)
                        if os.path.exists(old_file_path):
                            try:
                                os.remove(old_file_path)
                            except Exception as e:
                                print(f"Error deleting old profile picture: {e}")
                    
                    filename = secure_filename(f"profile_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                    os.makedirs(app.config['PROFILE_PICTURES_FOLDER'], exist_ok=True)
                    file_path = os.path.join(app.config['PROFILE_PICTURES_FOLDER'], filename)
                    file.save(file_path)
                    musician.profile_picture = f"profiles/pictures/{filename}"
            
            # Handle banner upload
            if form.banner.data:
                file = form.banner.data
                if hasattr(file, 'filename') and file.filename:
                    # Delete old banner if exists
                    if musician.banner:
                        old_file_path = os.path.join(app.static_folder, musician.banner)
                        if os.path.exists(old_file_path):
                            try:
                                os.remove(old_file_path)
                            except Exception as e:
                                print(f"Error deleting old banner: {e}")
                    
                    filename = secure_filename(f"banner_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                    os.makedirs(app.config['BANNERS_FOLDER'], exist_ok=True)
                    file_path = os.path.join(app.config['BANNERS_FOLDER'], filename)
                    file.save(file_path)
                    musician.banner = f"profiles/banners/{filename}"
            
            db.session.commit()
            flash('Profile picture and banner updated successfully.', 'success')
            return redirect(url_for('view_musician_profile', id=id))
    
    return render_template('edit_picture.html', form=form, musician=musician)


@app.route('/musicians/<int:id>/delete-profile-picture', methods=['POST'])
@login_required
def delete_profile_picture(id):
    """Delete profile picture"""
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Check permissions
    can_edit = (user and current_user.id == user.id) or current_user.is_worship_leader()
    if not can_edit:
        flash('You do not have permission to edit this profile.', 'danger')
        return redirect(url_for('view_musician_profile', id=id))
    
    if musician.profile_picture:
        # Delete file from disk
        file_path = os.path.join(app.static_folder, musician.profile_picture)
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as e:
                print(f"Error deleting profile picture: {e}")
        
        # Clear from database
        musician.profile_picture = None
        db.session.commit()
        flash('Profile picture deleted successfully.', 'success')
    else:
        flash('No profile picture to delete.', 'info')
    
    return redirect(url_for('edit_musician_picture', id=id))


@app.route('/musicians/<int:id>/delete-banner', methods=['POST'])
@login_required
def delete_banner(id):
    """Delete banner image"""
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Check permissions
    can_edit = (user and current_user.id == user.id) or current_user.is_worship_leader()
    if not can_edit:
        flash('You do not have permission to edit this profile.', 'danger')
        return redirect(url_for('view_musician_profile', id=id))
    
    if musician.banner:
        # Delete file from disk
        file_path = os.path.join(app.static_folder, musician.banner)
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as e:
                print(f"Error deleting banner: {e}")
        
        # Clear from database
        musician.banner = None
        db.session.commit()
        flash('Banner deleted successfully.', 'success')
    else:
        flash('No banner to delete.', 'info')
    
    return redirect(url_for('edit_musician_picture', id=id))


@app.route('/musicians/<int:id>/customize', methods=['GET', 'POST'])
@login_required
def customize_profile(id):
    """Friendster-like profile customization page"""
    musician = Musician.query.get_or_404(id)
    user = musician.user if musician.user_id else None
    
    # Check permissions - only profile owner can customize
    if not user or current_user.id != user.id:
        flash('You can only customize your own profile.', 'danger')
        return redirect(url_for('view_musician_profile', id=id))
    
    form = ProfileCustomizationForm(obj=musician)
    
    if form.validate_on_submit():
        # Clear any existing background image
        if musician.background_image:
            file_path = os.path.join(app.static_folder, musician.background_image)
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except Exception as e:
                    print(f"Error deleting background image: {e}")
            musician.background_image = None
        
        # Update customization fields
        musician.background_color = form.background_color.data or None
        musician.text_color = form.text_color.data or None
        musician.link_color = form.link_color.data or None
        musician.profile_theme = form.profile_theme.data or 'default'
        musician.music_player_embed = form.music_player_embed.data or None
        musician.custom_css = form.custom_css.data or None
        
        db.session.commit()
        flash('Profile customization saved successfully!', 'success')
        return redirect(url_for('view_musician_profile', id=id))
    
    # Pre-populate form with current values
    if request.method == 'GET':
        form.background_color.data = musician.background_color or ''
        form.text_color.data = musician.text_color or ''
        form.link_color.data = musician.link_color or ''
        form.profile_theme.data = musician.profile_theme or 'default'
        form.music_player_embed.data = musician.music_player_embed or ''
        form.custom_css.data = musician.custom_css or ''
    
    return render_template('customize_profile.html', form=form, musician=musician)


@app.route('/musicians/<int:id>/delete', methods=['POST'])
@worship_leader_required
def delete_musician(id):
    musician = Musician.query.get_or_404(id)
    db.session.delete(musician)
    db.session.commit()
    flash('Team member deleted successfully.', 'success')
    return redirect(url_for('musicians'))


# Post interaction routes
@app.route('/posts/<int:post_id>/edit', methods=['GET', 'POST'])
@login_required
def edit_profile_post(post_id):
    """Edit a profile post - only the post owner can edit"""
    post = ProfilePost.query.get_or_404(post_id)
    musician = post.musician
    user = musician.user if musician.user_id else None
    
    # Check permissions - only post owner can edit
    if not user or current_user.id != user.id:
        flash('You can only edit your own posts.', 'danger')
        return redirect(url_for('view_musician_profile', id=musician.id))
    
    form = ProfilePostForm(obj=post)
    if form.validate_on_submit():
        post.content = form.content.data or None
        
        # Handle image upload
        if form.image.data:
            file = form.image.data
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"post_img_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['POSTS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['POSTS_FOLDER'], filename)
                file.save(file_path)
                post.image_path = f"profiles/posts/{filename}"
        
        # Handle video upload
        if form.video.data:
            file = form.video.data
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"post_vid_{musician.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['POSTS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['POSTS_FOLDER'], filename)
                file.save(file_path)
                post.video_path = f"profiles/posts/{filename}"
        
        post.updated_at = datetime.utcnow()
        db.session.commit()
        flash('Post updated successfully!', 'success')
        return redirect(url_for('view_musician_profile', id=musician.id))
    
    return render_template('edit_post.html', form=form, post=post, musician=musician)


@app.route('/posts/<int:post_id>/delete', methods=['POST'])
@login_required
def delete_profile_post(post_id):
    """Delete a profile post - only the post owner can delete"""
    post = ProfilePost.query.get_or_404(post_id)
    musician = post.musician
    user = musician.user if musician.user_id else None
    
    # Check permissions - only post owner can delete
    if not user or current_user.id != user.id:
        flash('You can only delete your own posts.', 'danger')
        return redirect(url_for('view_musician_profile', id=musician.id))
    
    # Delete associated media files
    if post.image_path:
        image_path = os.path.join(app.root_path, 'static', post.image_path)
        if os.path.exists(image_path):
            try:
                os.remove(image_path)
            except:
                pass
    
    if post.video_path:
        video_path = os.path.join(app.root_path, 'static', post.video_path)
        if os.path.exists(video_path):
            try:
                os.remove(video_path)
            except:
                pass
    
    db.session.delete(post)
    db.session.commit()
    flash('Post deleted successfully.', 'success')
    return redirect(url_for('view_musician_profile', id=musician.id))


@app.route('/posts/<int:post_id>/like', methods=['POST'])
@csrf.exempt
@login_required
def toggle_post_like(post_id):
    """Toggle like on a post - removes heart if exists"""
    post = ProfilePost.query.get_or_404(post_id)
    existing_like = PostLike.query.filter_by(post_id=post_id, user_id=current_user.id).first()
    existing_heart = PostHeart.query.filter_by(post_id=post_id, user_id=current_user.id).first()
    
    if existing_like:
        # Remove like
        db.session.delete(existing_like)
        action = 'unliked'
    else:
        # Remove heart if exists (can only have one reaction)
        if existing_heart:
            db.session.delete(existing_heart)
        # Add like
        like = PostLike(post_id=post_id, user_id=current_user.id)
        db.session.add(like)
        action = 'liked'
        
        # Create notification for post owner (if not liking own post)
        if post.musician and post.musician.user_id and post.musician.user_id != current_user.id:
            notification = Notification(
                user_id=post.musician.user_id,
                notification_type='like',
                actor_id=current_user.id,
                post_id=post_id
            )
            db.session.add(notification)
    
    db.session.commit()
    # Refresh post to get updated counts
    db.session.refresh(post)
    return jsonify({
        'success': True, 
        'action': action, 
        'like_count': len(post.likes),
        'heart_count': len(post.hearts),
        'has_like': post.is_liked_by(current_user.id),
        'has_heart': post.is_hearted_by(current_user.id)
    })


@app.route('/posts/<int:post_id>/heart', methods=['POST'])
@csrf.exempt
@login_required
def toggle_post_heart(post_id):
    """Toggle heart on a post - removes like if exists"""
    post = ProfilePost.query.get_or_404(post_id)
    existing_heart = PostHeart.query.filter_by(post_id=post_id, user_id=current_user.id).first()
    existing_like = PostLike.query.filter_by(post_id=post_id, user_id=current_user.id).first()
    
    if existing_heart:
        # Remove heart
        db.session.delete(existing_heart)
        action = 'unhearted'
    else:
        # Remove like if exists (can only have one reaction)
        if existing_like:
            db.session.delete(existing_like)
        # Add heart
        heart = PostHeart(post_id=post_id, user_id=current_user.id)
        db.session.add(heart)
        action = 'hearted'
        
        # Create notification for post owner (if not hearting own post)
        if post.musician and post.musician.user_id and post.musician.user_id != current_user.id:
            notification = Notification(
                user_id=post.musician.user_id,
                notification_type='heart',
                actor_id=current_user.id,
                post_id=post_id
            )
            db.session.add(notification)
    
    db.session.commit()
    # Refresh post to get updated counts
    db.session.refresh(post)
    return jsonify({
        'success': True, 
        'action': action, 
        'like_count': len(post.likes),
        'heart_count': len(post.hearts),
        'has_like': post.is_liked_by(current_user.id),
        'has_heart': post.is_hearted_by(current_user.id)
    })


@app.route('/posts/<int:post_id>/share', methods=['POST'])
@login_required
def share_post(post_id):
    """Share a post to the current user's wall with editable content"""
    original_post = ProfilePost.query.get_or_404(post_id)
    
    # Get or create musician profile for current user
    if not current_user.musician:
        musician = Musician(
            name=current_user.get_display_name(),
            user_id=current_user.id
        )
        db.session.add(musician)
        db.session.commit()  # Commit to ensure musician.id is available
    else:
        musician = current_user.musician
    
    # Get content from form (user can edit it)
    content = request.form.get('content', '').strip()
    if not content:
        flash('Please add some content before sharing.', 'warning')
        return redirect(url_for('view_musician_profile', id=original_post.musician_id))
    
    # Check if user already shared this post
    existing_share = PostRepost.query.filter_by(post_id=post_id, user_id=current_user.id).first()
    if existing_share:
        flash('You have already shared this post.', 'warning')
        return redirect(url_for('view_musician_profile', id=original_post.musician_id))
    
    # Create a new post on the sharer's wall
    shared_post = ProfilePost(musician_id=musician.id)
    shared_post.content = content
    
    # Copy image if exists (from original post)
    original_image_path = request.form.get('original_image_path')
    if original_image_path:
        shared_post.image_path = original_image_path
    
    # Copy video if exists (from original post)
    original_video_path = request.form.get('original_video_path')
    if original_video_path:
        shared_post.video_path = original_video_path
    
    db.session.add(shared_post)
    
    # Record the share
    share_record = PostRepost(post_id=post_id, user_id=current_user.id)
    db.session.add(share_record)
    
    # Create notification for original post owner (if not sharing own post)
    if original_post.musician and original_post.musician.user_id and original_post.musician.user_id != current_user.id:
        notification = Notification(
            user_id=original_post.musician.user_id,
            notification_type='share',
            actor_id=current_user.id,
            post_id=post_id
        )
        db.session.add(notification)
    
    try:
        db.session.commit()
        # Verify the post was created and refresh to get the ID
        db.session.refresh(shared_post)
        
        if shared_post.id and shared_post.musician_id == musician.id:
            flash('Post shared to your wall!', 'success')
            # Redirect to the sharer's own profile to see the shared post
            # Add a timestamp to force refresh
            return redirect(url_for('view_musician_profile', id=musician.id) + '?shared=1')
        else:
            flash('Error: Post was not created correctly. Please try again.', 'danger')
            return redirect(url_for('view_musician_profile', id=original_post.musician_id))
    except Exception as e:
        db.session.rollback()
        flash(f'Error sharing post: {str(e)}', 'danger')
        return redirect(url_for('view_musician_profile', id=original_post.musician_id))


@app.route('/posts/<int:post_id>/comment', methods=['POST'])
@login_required
def add_post_comment(post_id):
    """Add a comment to a post"""
    post = ProfilePost.query.get_or_404(post_id)
    content = request.form.get('content', '').strip()
    
    if not content:
        flash('Comment cannot be empty.', 'warning')
        return redirect(url_for('view_musician_profile', id=post.musician_id))
    
    comment = PostComment(
        post_id=post_id,
        user_id=current_user.id,
        content=content
    )
    db.session.add(comment)
    db.session.flush()  # Flush to get comment.id
    
    # Create notification for post owner (if not commenting on own post)
    if post.musician and post.musician.user_id and post.musician.user_id != current_user.id:
        notification = Notification(
            user_id=post.musician.user_id,
            notification_type='comment',
            actor_id=current_user.id,
            post_id=post_id,
            comment_id=comment.id
        )
        db.session.add(notification)
    
    db.session.commit()
    flash('Comment added successfully.', 'success')
    return redirect(url_for('view_musician_profile', id=post.musician_id))


@app.route('/comments/<int:comment_id>/delete', methods=['POST'])
@login_required
def delete_post_comment(comment_id):
    """Delete a comment - only comment owner can delete"""
    comment = PostComment.query.get_or_404(comment_id)
    post = comment.post
    
    # Check permissions - only comment owner can delete
    if comment.user_id != current_user.id:
        flash('You can only delete your own comments.', 'danger')
        return redirect(url_for('view_musician_profile', id=post.musician_id))
    
    db.session.delete(comment)
    db.session.commit()
    flash('Comment deleted successfully.', 'success')
    return redirect(url_for('view_musician_profile', id=post.musician_id))


# Musician Availability routes - now works with user_id
@app.route('/users/<int:user_id>/availability')
@login_required
def musician_availability(user_id):
    user = User.query.get_or_404(user_id)
    
    # Allow viewing if user is viewing their own or if current user is worship leader
    # Regular users can only view their own availability
    
    # Get or create musician profile for this user
    musician = user.musician
    if not musician:
        # Auto-create musician profile if it doesn't exist
        musician = Musician(
            name=user.get_display_name(),
            user_id=user.id
        )
        db.session.add(musician)
        db.session.commit()
    
    # Get all availability records for this musician (only show approved leaves)
    availability_records = MusicianAvailability.query.filter_by(
        musician_id=musician.id
    ).all()
    
    # Also check for approved leave requests and add them to availability
    approved_leaves = LeaveRequest.query.filter_by(
        musician_id=musician.id,
        status='approved'
    ).all()
    
    # Convert to dictionary for easy lookup
    availability_dict = {}
    
    # Add availability records
    for record in availability_records:
        availability_dict[record.date.isoformat()] = {
            'is_available': record.is_available,
            'notes': record.notes
        }
    
    # Add approved leave requests (these override availability records)
    for leave in approved_leaves:
        availability_dict[leave.date.isoformat()] = {
            'is_available': False,
            'notes': leave.reason,
            'is_pending': False,  # Explicitly mark as not pending
            'is_approved': True  # Mark as approved leave
        }
    
    # Check for pending leave requests - show them if user is viewing their own availability
    # Only add if there's no approved leave for that date
    pending_leaves = []
    if user.id == current_user.id:
        pending_leaves = LeaveRequest.query.filter_by(
            musician_id=musician.id,
            status='pending'
        ).all()
        # Add pending leaves to availability dict with special marker
        # Only if there's no approved leave for that date (approved takes precedence)
        for leave in pending_leaves:
            date_key = leave.date.isoformat()
            # Only add if this date doesn't exist or if existing entry is not an approved leave
            # Approved leaves have is_pending=False (or not set)
            existing = availability_dict.get(date_key, {})
            # Don't overwrite if there's already an approved leave (is_pending=False)
            if date_key not in availability_dict or existing.get('is_pending') != False:
                availability_dict[date_key] = {
                    'is_available': False,
                    'notes': leave.reason,
                    'is_pending': True,
                    'leave_request_id': leave.id
                }
    
    # Also check for rejected leave requests to show them (but not as unavailable)
    rejected_leaves = LeaveRequest.query.filter_by(
        musician_id=musician.id,
        status='rejected'
    ).all()
    
    # Note: Rejected leaves are not added to availability_dict, so they won't show in calendar
    # They can be viewed in the leave requests page
    
    # Check if user is viewing their own availability
    is_own_availability = (user.id == current_user.id)
    
    return render_template('musician_availability.html', 
                         musician=musician,
                         user=user,  # Always pass user
                         availability=availability_dict,
                         pending_leaves=pending_leaves,
                         is_own_availability=is_own_availability)


@app.route('/users/<int:user_id>/availability/toggle', methods=['POST'])
@csrf.exempt
@login_required
def toggle_availability(user_id):
    """File a leave request (replaces direct availability toggle)"""
    try:
        user = User.query.get_or_404(user_id)
        
        # Only allow users to file leave for themselves
        if user.id != current_user.id:
            return jsonify({'success': False, 'message': 'You can only file a leave for yourself.'}), 403
        
        # Get or create team member profile for this user
        musician = user.musician
        if not musician:
            musician = Musician(
                name=user.get_display_name(),
                user_id=user.id
            )
            db.session.add(musician)
            db.session.flush()  # Flush to get ID without committing
        
        musician_id = musician.id
        if not musician_id:
            return jsonify({'success': False, 'message': 'Failed to create team member profile.'}), 500
        
        data = request.get_json()
        if not data:
            return jsonify({'success': False, 'message': 'Invalid request data'}), 400
        
        date_str = data.get('date')
        is_available = data.get('is_available', True)
        reason = data.get('reason', '') or data.get('notes', '')
        
        # Strip whitespace and validate
        if reason:
            reason = reason.strip()
        
        if not date_str:
            return jsonify({'success': False, 'message': 'Date is required'}), 400
        
        if not is_available and not reason:
            return jsonify({'success': False, 'message': 'Reason is required for leave requests.'}), 400
        
        try:
            from datetime import datetime
            date_obj = datetime.strptime(date_str, '%Y-%m-%d').date()
        except ValueError:
            return jsonify({'success': False, 'message': 'Invalid date format'}), 400
        
        # If setting as available, directly update availability
        if is_available:
            availability = MusicianAvailability.query.filter_by(
                musician_id=musician_id,
                date=date_obj
            ).first()
            
            if availability:
                availability.is_available = True
                availability.notes = ''
            else:
                availability = MusicianAvailability(
                    musician_id=musician_id,
                    date=date_obj,
                    is_available=True,
                    notes=''
                )
                db.session.add(availability)
            db.session.commit()
            return jsonify({'success': True, 'message': 'Availability updated'})
        
        # For leave requests (unavailable), create a leave request
        # Check if leave request already exists for this date (pending or approved)
        existing_request = LeaveRequest.query.filter_by(
            user_id=user_id,
            musician_id=musician_id,
            date=date_obj
        ).filter(LeaveRequest.status.in_(['pending', 'approved'])).first()
        
        if existing_request:
            return jsonify({'success': False, 'message': 'You already have a leave request for this date.'}), 400
        
        # Ensure reason is not empty
        if not reason or not reason.strip():
            return jsonify({'success': False, 'message': 'A reason is required for leave requests.'}), 400
        
        # Check if user is a team leader - auto-approve if so
        is_team_leader = current_user.is_team_leader()
        
        # Create leave request with appropriate status
        leave_request = LeaveRequest(
            user_id=user_id,
            musician_id=musician_id,
            date=date_obj,
            reason=reason.strip(),
            status='approved' if is_team_leader else 'pending'
        )
        
        # If team leader, auto-approve by setting reviewed fields
        if is_team_leader:
            leave_request.reviewed_by = current_user.id
            leave_request.reviewed_at = datetime.utcnow()
        
        db.session.add(leave_request)
        
        # If auto-approved (team leader), create availability record immediately
        if is_team_leader:
            availability = MusicianAvailability.query.filter_by(
                musician_id=musician_id,
                date=date_obj
            ).first()
            
            if availability:
                availability.is_available = False
                availability.notes = reason.strip()
            else:
                availability = MusicianAvailability(
                    musician_id=musician_id,
                    date=date_obj,
                    is_available=False,
                    notes=reason.strip()
                )
                db.session.add(availability)
        
        # Commit leave request and availability (if team leader)
        try:
            db.session.commit()
        except Exception as commit_error:
            db.session.rollback()
            error_str = str(commit_error).lower()
            if 'leave_request' in error_str or 'no such table' in error_str:
                return jsonify({
                    'success': False, 
                    'message': 'Database migration required. Please restart the application to run migrations.'
                }), 500
            raise
        
        # Log activity
        if is_team_leader:
            log_activity(
                activity_type='leave_approved',
                actor_id=current_user.id,
                target_user_id=user_id,
                description=f"{current_user.get_display_name()} filed and auto-approved a leave request for {date_obj.strftime('%B %d, %Y')}: {reason.strip()}",
                leave_request_id=leave_request.id,
                metadata={'date': date_obj.isoformat(), 'reason': reason.strip()}
            )
        else:
            log_activity(
                activity_type='leave_filed',
                actor_id=current_user.id,
                target_user_id=user_id,
                description=f"{current_user.get_display_name()} filed a leave request for {date_obj.strftime('%B %d, %Y')}: {reason.strip()}",
                leave_request_id=leave_request.id,
                metadata={'date': date_obj.isoformat(), 'reason': reason.strip()}
            )
        
        # Only send notifications to other team leaders if not auto-approved
        if not is_team_leader:
            team_leaders = User.query.filter(User.role.in_(['team_leader', 'admin'])).all()
            for team_leader in team_leaders:
                notification = Notification(
                    user_id=team_leader.id,
                    notification_type='leave_request',
                    actor_id=user_id,
                    is_read=False
                )
                # Set leave_request_id if the attribute exists (column may not exist in old databases)
                if hasattr(Notification, 'leave_request_id'):
                    try:
                        notification.leave_request_id = leave_request.id
                    except Exception:
                        # Column doesn't exist yet - continue without it
                        pass
                db.session.add(notification)
            
            try:
                db.session.commit()
            except Exception as commit_error:
                # If notification commit fails due to missing column, still return success
                # The leave request was already created
                error_str = str(commit_error).lower()
                if 'leave_request_id' in error_str or 'no such column' in error_str:
                    print(f"Warning: Could not set leave_request_id on notifications: {commit_error}")
                    db.session.rollback()
                    # Re-commit without leave_request_id
                    for team_leader in team_leaders:
                        notification = Notification(
                            user_id=team_leader.id,
                            notification_type='leave_request',
                            actor_id=user_id,
                            is_read=False
                        )
                        db.session.add(notification)
                    db.session.commit()
                else:
                    raise
        
        # Return appropriate message based on approval status
        if is_team_leader:
            return jsonify({
                'success': True,
                'message': 'Leave request approved automatically. It has been added to your schedule.'
            })
        else:
            return jsonify({
                'success': True,
                'message': 'Leave request submitted successfully. It will appear in your schedule once approved by your Team Leader.'
            })
    
    except Exception as e:
        db.session.rollback()
        import traceback
        error_traceback = traceback.format_exc()
        print(f"Error filing leave request: {str(e)}")
        print(error_traceback)
        
        # Provide more helpful error messages
        error_message = str(e)
        if 'no such column' in error_message.lower() or 'leave_request_id' in error_message.lower():
            error_message = 'Database migration required. Please restart the application to run migrations.'
        elif 'no such table' in error_message.lower() or 'leave_request' in error_message.lower():
            error_message = 'Database migration required. Please restart the application to run migrations.'
        
        return jsonify({'success': False, 'message': f'An error occurred: {error_message}'}), 500


@app.route('/users/<int:user_id>/availability/remove', methods=['POST'])
@csrf.exempt
@login_required
def remove_availability(user_id):
    user = User.query.get_or_404(user_id)
    
    # Check if user has permission
    if not current_user.is_worship_leader() and user.id != current_user.id:
        return jsonify({'success': False, 'message': 'You can only remove your own availability.'}), 403
    
        # Get team member profile
    musician = user.musician
    if not musician:
        return jsonify({'success': False, 'message': 'Team member profile not found.'}), 404
    
    musician_id = musician.id
    
    data = request.get_json()
    date_str = data.get('date')
    
    if not date_str:
        return jsonify({'success': False, 'message': 'Date is required'}), 400
    
    try:
        from datetime import datetime
        date_obj = datetime.strptime(date_str, '%Y-%m-%d').date()
    except ValueError:
        return jsonify({'success': False, 'message': 'Invalid date format'}), 400
    
    # First check if it's an approved leave request
    leave_request = LeaveRequest.query.filter_by(
        musician_id=musician_id,
        date=date_obj,
        status='approved'
    ).first()
    
    if leave_request:
        # Cancel the approved leave request
        # Only the user who requested it can cancel it
        if leave_request.user_id != current_user.id:
            return jsonify({'success': False, 'message': 'You can only cancel your own leave requests.'}), 403
        
        # Update leave request status to cancelled
        leave_request.status = 'cancelled'
        leave_request.reviewed_at = datetime.utcnow()
        
        # Delete the associated availability record
        availability = MusicianAvailability.query.filter_by(
            musician_id=musician_id,
            date=date_obj
        ).first()
        if availability:
            db.session.delete(availability)
        
        # Notify the approver (if there was one) that the leave was cancelled
        if leave_request.reviewed_by:
            notification = Notification(
                user_id=leave_request.reviewed_by,
                notification_type='leave_cancelled',
                actor_id=current_user.id,
                is_read=False
            )
            # Set leave_request_id if the attribute exists
            if hasattr(Notification, 'leave_request_id'):
                try:
                    notification.leave_request_id = leave_request.id
                except Exception:
                    pass
            db.session.add(notification)
        
        # Log activity
        log_activity(
            activity_type='leave_cancelled',
            actor_id=current_user.id,
            target_user_id=leave_request.user_id,
            description=f"{current_user.get_display_name()} cancelled their approved leave request for {leave_request.date.strftime('%B %d, %Y')}: {leave_request.reason}",
            leave_request_id=leave_request.id,
            metadata={'date': leave_request.date.isoformat(), 'reason': leave_request.reason}
        )
        
        try:
            db.session.commit()
            return jsonify({'success': True, 'message': 'Leave request cancelled successfully. The approver has been notified.'})
        except Exception as e:
            db.session.rollback()
            return jsonify({'success': False, 'message': f'Error cancelling leave request: {str(e)}'}), 500
    
    # Delete the availability record
    availability = MusicianAvailability.query.filter_by(
        musician_id=musician_id,
        date=date_obj
    ).first()
    
    if availability:
        db.session.delete(availability)
        db.session.commit()
        return jsonify({'success': True, 'message': 'Availability removed'})
    else:
        # Record doesn't exist - might have already been deleted or never existed
        # Return success anyway to avoid error popup
        return jsonify({'success': True, 'message': 'Availability record not found or already removed'})


@app.route('/leave-requests')
@login_required
def leave_requests():
    """View leave requests - Team Leaders see all pending, users see their own"""
    if current_user.is_team_leader():
        # Team Leaders see all pending requests
        pending_requests = LeaveRequest.query.filter_by(status='pending').order_by(LeaveRequest.date.asc()).all()
        return render_template('leave_requests.html', requests=pending_requests, is_team_leader=True)
    else:
        # Regular users see their own requests
        user_requests = LeaveRequest.query.filter_by(user_id=current_user.id).order_by(LeaveRequest.date.desc()).all()
        return render_template('leave_requests.html', requests=user_requests, is_team_leader=False)


@app.route('/leave-requests/<int:request_id>/approve', methods=['POST'])
@csrf.exempt
@login_required
def approve_leave_request(request_id):
    """Approve a leave request - Team Leaders only"""
    if not current_user.is_team_leader():
        return jsonify({'success': False, 'message': 'Only Team Leaders can approve leave requests.'}), 403
    
    leave_request = LeaveRequest.query.get_or_404(request_id)
    
    if leave_request.status != 'pending':
        return jsonify({'success': False, 'message': 'This leave request has already been processed.'}), 400
    
    # Update leave request status
    leave_request.status = 'approved'
    leave_request.reviewed_by = current_user.id
    leave_request.reviewed_at = datetime.utcnow()
    
    # Create or update availability record
    availability = MusicianAvailability.query.filter_by(
        musician_id=leave_request.musician_id,
        date=leave_request.date
    ).first()
    
    if availability:
        availability.is_available = False
        availability.notes = leave_request.reason
    else:
        availability = MusicianAvailability(
            musician_id=leave_request.musician_id,
            date=leave_request.date,
            is_available=False,
            notes=leave_request.reason
        )
        db.session.add(availability)
    
    # Notify the user that their leave was approved
    notification = Notification(
        user_id=leave_request.user_id,
        notification_type='leave_approved',
        actor_id=current_user.id,
        is_read=False
    )
    # Set leave_request_id if the attribute exists (column may not exist in old databases)
    if hasattr(Notification, 'leave_request_id'):
        try:
            notification.leave_request_id = leave_request.id
        except Exception:
            pass  # Ignore if column doesn't exist
    db.session.add(notification)
    
    # Log activity
    log_activity(
        activity_type='leave_approved',
        actor_id=current_user.id,
        target_user_id=leave_request.user_id,
        description=f"{current_user.get_display_name()} approved {leave_request.user.get_display_name()}'s leave request for {leave_request.date.strftime('%B %d, %Y')}: {leave_request.reason}",
        leave_request_id=leave_request.id,
        metadata={'date': leave_request.date.isoformat(), 'reason': leave_request.reason}
    )
    
    try:
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error approving leave request: {str(e)}'}), 500
    
    return jsonify({'success': True, 'message': 'Leave request approved successfully.'})


@app.route('/leave-requests/<int:request_id>/reject', methods=['POST'])
@csrf.exempt
@login_required
def reject_leave_request(request_id):
    """Reject a leave request - Team Leaders only"""
    if not current_user.is_team_leader():
        return jsonify({'success': False, 'message': 'Only Team Leaders can reject leave requests.'}), 403
    
    leave_request = LeaveRequest.query.get_or_404(request_id)
    
    if leave_request.status != 'pending':
        return jsonify({'success': False, 'message': 'This leave request has already been processed.'}), 400
    
    data = request.get_json()
    review_notes = data.get('review_notes', '') if data else ''
    
    # Update leave request status
    leave_request.status = 'rejected'
    leave_request.reviewed_by = current_user.id
    leave_request.reviewed_at = datetime.utcnow()
    leave_request.review_notes = review_notes
    
    # Notify the user that their leave was rejected
    notification = Notification(
        user_id=leave_request.user_id,
        notification_type='leave_rejected',
        actor_id=current_user.id,
        is_read=False
    )
    # Set leave_request_id if the attribute exists (column may not exist in old databases)
    if hasattr(Notification, 'leave_request_id'):
        try:
            notification.leave_request_id = leave_request.id
        except Exception:
            pass  # Ignore if column doesn't exist
    db.session.add(notification)
    
    try:
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error rejecting leave request: {str(e)}'}), 500
    
    return jsonify({'success': True, 'message': 'Leave request rejected.'})


@app.route('/leave-requests/approve-all', methods=['POST'])
@csrf.exempt
@login_required
def approve_all_leave_requests():
    """Approve all pending leave requests - Team Leaders only"""
    if not current_user.is_team_leader():
        return jsonify({'success': False, 'message': 'Only Team Leaders can approve leave requests.'}), 403
    
    try:
        # Get all pending leave requests
        pending_requests = LeaveRequest.query.filter_by(status='pending').all()
        
        if not pending_requests:
            return jsonify({'success': False, 'message': 'No pending leave requests to approve.'}), 400
        
        approved_count = 0
        errors = []
        
        for leave_request in pending_requests:
            try:
                # Update leave request status
                leave_request.status = 'approved'
                leave_request.reviewed_by = current_user.id
                leave_request.reviewed_at = datetime.utcnow()
                
                # Create or update availability record
                availability = MusicianAvailability.query.filter_by(
                    musician_id=leave_request.musician_id,
                    date=leave_request.date
                ).first()
                
                if availability:
                    availability.is_available = False
                    availability.notes = leave_request.reason
                else:
                    availability = MusicianAvailability(
                        musician_id=leave_request.musician_id,
                        date=leave_request.date,
                        is_available=False,
                        notes=leave_request.reason
                    )
                    db.session.add(availability)
                
                # Notify the user that their leave was approved
                notification = Notification(
                    user_id=leave_request.user_id,
                    notification_type='leave_approved',
                    actor_id=current_user.id,
                    is_read=False
                )
                # Set leave_request_id if the attribute exists
                if hasattr(Notification, 'leave_request_id'):
                    try:
                        notification.leave_request_id = leave_request.id
                    except Exception:
                        pass
                db.session.add(notification)
                
                # Log activity for each approved leave
                log_activity(
                    activity_type='leave_approved',
                    actor_id=current_user.id,
                    target_user_id=leave_request.user_id,
                    description=f"{current_user.get_display_name()} approved {leave_request.user.get_display_name()}'s leave request for {leave_request.date.strftime('%B %d, %Y')}: {leave_request.reason}",
                    leave_request_id=leave_request.id,
                    metadata={'date': leave_request.date.isoformat(), 'reason': leave_request.reason}
                )
                
                approved_count += 1
            except Exception as e:
                errors.append(f"Error approving request {leave_request.id}: {str(e)}")
                db.session.rollback()
                continue
        
        # Commit all changes
        try:
            db.session.commit()
        except Exception as e:
            db.session.rollback()
            return jsonify({'success': False, 'message': f'Error committing changes: {str(e)}'}), 500
        
        if errors:
            message = f'Approved {approved_count} leave request(s). Some errors occurred: {"; ".join(errors)}'
        else:
            message = f'Successfully approved {approved_count} leave request(s).'
        
        return jsonify({'success': True, 'message': message, 'count': approved_count})
    
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error approving leave requests: {str(e)}'}), 500


# Sunday Service routes
@app.route('/services')
@login_required
def services():
    services_list = SundayService.query.order_by(SundayService.date.desc()).all()
    return render_template('services.html', services=services_list)


@app.route('/services/add', methods=['GET', 'POST'])
@worship_leader_required
def add_service():
    form = ServiceForm()
    if form.validate_on_submit():
        service = SundayService(
            date=form.date.data,
            theme=form.theme.data or None,
            notes=form.notes.data or None,
            created_by=current_user.id
        )
        db.session.add(service)
        db.session.commit()
        flash('Service added successfully.', 'success')
        return redirect(url_for('service_detail', id=service.id))
    return render_template('service_form.html', form=form, title='Add Sunday Service')


@app.route('/services/<int:id>')
@login_required
def service_detail(id):
    service = SundayService.query.get_or_404(id)
    musicians_list = Musician.query.order_by(Musician.name).all()
    form = ServiceMusicianForm()
    form.musician_id.choices = [(m.id, m.name) for m in musicians_list]
    return render_template('service_detail.html', service=service, musicians=musicians_list, form=form)


@app.route('/services/<int:id>/edit', methods=['GET', 'POST'])
@worship_leader_required
def edit_service(id):
    service = SundayService.query.get_or_404(id)
    form = ServiceForm(obj=service)
    if form.validate_on_submit():
        service.date = form.date.data
        service.theme = form.theme.data or None
        service.notes = form.notes.data or None
        db.session.commit()
        flash('Service updated successfully.', 'success')
        return redirect(url_for('service_detail', id=service.id))
    return render_template('service_form.html', form=form, title='Edit Sunday Service', service=service)


@app.route('/services/<int:id>/delete', methods=['POST'])
@worship_leader_required
def delete_service(id):
    service = SundayService.query.get_or_404(id)
    db.session.delete(service)
    db.session.commit()
    flash('Service deleted successfully.', 'success')
    return redirect(url_for('services'))


@app.route('/services/<int:service_id>/musicians/add', methods=['POST'])
@worship_leader_required
def add_service_musician(service_id):
    service = SundayService.query.get_or_404(service_id)
    form = ServiceMusicianForm()
    form.musician_id.choices = [(m.id, m.name) for m in Musician.query.order_by(Musician.name).all()]
    
    if form.validate_on_submit():
        assignment = ServiceMusician(
            service_id=service_id,
            musician_id=form.musician_id.data,
            instrument=form.instrument.data,
            role=form.role.data or None
        )
        db.session.add(assignment)
        db.session.commit()
        flash('Team member added to service.', 'success')
    else:
        flash('Error adding team member.', 'danger')
    return redirect(url_for('service_detail', id=service_id))


@app.route('/services/<int:service_id>/musicians/<int:assignment_id>/delete', methods=['POST'])
@worship_leader_required
def delete_service_musician(service_id, assignment_id):
    assignment = ServiceMusician.query.get_or_404(assignment_id)
    db.session.delete(assignment)
    db.session.commit()
    flash('Team member removed from service.', 'success')
    return redirect(url_for('service_detail', id=service_id))


# Practice routes
@app.route('/practices')
@login_required
def practices():
    practices_list = Practice.query.order_by(Practice.date.desc()).all()
    return render_template('practices.html', practices=practices_list)


@app.route('/practices/add', methods=['GET', 'POST'])
@worship_leader_required
def add_practice():
    from datetime import datetime
    form = PracticeForm()
    if form.validate_on_submit():
        # Convert string date to date object
        date_obj = None
        if form.date.data:
            if isinstance(form.date.data, str):
                date_obj = datetime.strptime(form.date.data, '%Y-%m-%d').date()
            else:
                date_obj = form.date.data
        
        # Convert string time to time object
        time_obj = None
        if form.time.data:
            if isinstance(form.time.data, str):
                try:
                    # Try parsing as HH:MM:SS format first
                    time_obj = datetime.strptime(form.time.data, '%H:%M:%S').time()
                except ValueError:
                    # Fall back to HH:MM format
                    try:
                        time_obj = datetime.strptime(form.time.data, '%H:%M').time()
                    except ValueError:
                        time_obj = None
            else:
                time_obj = form.time.data
        
        practice = Practice(
            date=date_obj,
            time=time_obj,
            location=form.location.data or None,
            purpose=form.purpose.data or None,
            notes=form.notes.data or None,
            created_by=current_user.id
        )
        db.session.add(practice)
        db.session.flush()  # Flush to get practice.id
        
        # Create notifications for all users (except the creator)
        all_users = User.query.filter(User.id != current_user.id).all()
        for user in all_users:
            notification = Notification(
                user_id=user.id,
                notification_type='practice',
                actor_id=current_user.id,
                practice_id=practice.id
            )
            db.session.add(notification)
        
        db.session.commit()
        flash('Practice added successfully.', 'success')
        return redirect(url_for('practice_detail', id=practice.id))
    return render_template('practice_form.html', form=form, title='Schedule a Practice')


@app.route('/practices/<int:id>')
@login_required
def practice_detail(id):
    from sqlalchemy.orm import joinedload  # type: ignore
    from models import PracticeMusician, Musician
    # Load practice with all relationships
    practice = Practice.query.options(
        joinedload(Practice.musicians).joinedload(PracticeMusician.musician)
    ).get_or_404(id)
    
    # Ensure all musicians are properly loaded (fallback if joinedload doesn't work)
    # Query PracticeMusician directly to ensure we get all assignments
    practice_musicians = PracticeMusician.query.filter_by(practice_id=id).all()
    # Verify each assignment has its musician loaded
    for pm in practice_musicians:
        if pm.musician_id and not hasattr(pm, '_musician_loaded'):
            if not pm.musician:
                pm.musician = Musician.query.get(pm.musician_id)
            pm._musician_loaded = True
    
    # Update practice.musicians to use the directly queried list
    # This ensures we have all assignments with properly loaded musicians
    practice._musicians_list = practice_musicians
    musicians_list = Musician.query.order_by(Musician.name).all()
    songs_list = Song.query.order_by(Song.title).all()
    users_list = User.query.order_by(User.username).all()  # For "Prepared by" dropdown
    form = PracticeMusicianForm()
    # Show only users from the user list in dropdown
    # Get or create musician profile for each user
    musician_choices = []
    for user in users_list:
        if not user.musician:
            # Create musician profile if it doesn't exist
            musician = Musician(
                name=user.get_display_name(),
                user_id=user.id
            )
            db.session.add(musician)
            db.session.commit()
        else:
            musician = user.musician
        musician_choices.append((musician.id, user.get_display_name()))
    form.musician_id.choices = musician_choices
    
    # Get songs in lineup order
    practice_songs = PracticeSong.query.filter_by(practice_id=id).order_by(PracticeSong.order).all()
    
    return render_template('practice_detail.html', 
                         practice=practice, 
                         musicians=musicians_list, 
                         songs=songs_list,
                         users=users_list,
                         practice_songs=practice_songs,
                         form=form)


@app.route('/practices/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit_practice(id):
    from datetime import datetime
    practice = Practice.query.get_or_404(id)
    
    # Check permissions - only admin or worship leader can edit
    if not (current_user.is_admin() or current_user.is_worship_leader()):
        return redirect(url_for('practice_detail', id=id))
    
    form = PracticeForm(obj=practice)
    
    # Format date and time for HTML5 date/time inputs
    if practice.date:
        form.date.data = practice.date.strftime('%Y-%m-%d')
    if practice.time:
        form.time.data = practice.time.strftime('%H:%M')
    
    if form.validate_on_submit():
        # Convert string date to date object
        if form.date.data:
            if isinstance(form.date.data, str):
                practice.date = datetime.strptime(form.date.data, '%Y-%m-%d').date()
            else:
                practice.date = form.date.data
        
        # Convert string time to time object
        if form.time.data:
            if isinstance(form.time.data, str):
                try:
                    # Try parsing as HH:MM:SS format first
                    practice.time = datetime.strptime(form.time.data, '%H:%M:%S').time()
                except ValueError:
                    # Fall back to HH:MM format
                    try:
                        practice.time = datetime.strptime(form.time.data, '%H:%M').time()
                    except ValueError:
                        practice.time = None
            else:
                practice.time = form.time.data
        else:
            practice.time = None
        
        practice.location = form.location.data or None
        practice.purpose = form.purpose.data or None
        practice.notes = form.notes.data or None
        db.session.commit()
        flash('Practice updated successfully.', 'success')
        return redirect(url_for('practice_detail', id=practice.id))
    return render_template('practice_form.html', form=form, title='Edit Practice', practice=practice)


@app.route('/practices/<int:id>/delete', methods=['POST'])
@worship_leader_required
def delete_practice(id):
    practice = Practice.query.get_or_404(id)
    db.session.delete(practice)
    db.session.commit()
    flash('Practice deleted successfully.', 'success')
    return redirect(url_for('practices'))


@app.route('/practices/delete-all', methods=['POST'])
@login_required
@admin_required
def delete_all_practices():
    """Delete all scheduled practices - Admin only"""
    try:
        # Get count before deletion
        practice_count = Practice.query.count()
        
        # Delete all practices (cascade will handle related records)
        Practice.query.delete()
        db.session.commit()
        
        flash(f'Successfully deleted {practice_count} practice(s).', 'success')
    except Exception as e:
        db.session.rollback()
        flash(f'Error deleting practices: {str(e)}', 'danger')
    
    return redirect(url_for('practices'))


@app.route('/practices/<int:practice_id>/musicians/add', methods=['POST'])
@worship_leader_required
def add_practice_musician(practice_id):
    practice = Practice.query.get_or_404(practice_id)
    form = PracticeMusicianForm()
    # Show only users from the user list in dropdown
    # Get or create musician profile for each user
    users_list = User.query.order_by(User.username).all()
    musician_choices = []
    for user in users_list:
        if not user.musician:
            # Create musician profile if it doesn't exist, syncing role
            musician = Musician(
                name=user.get_display_name(),
                user_id=user.id,
                instruments=user.role if user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader'] else None
            )
            db.session.add(musician)
            db.session.commit()
        else:
            musician = user.musician
            # Sync role if musician profile exists but doesn't have a role set
            if not musician.instruments and user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
                musician.instruments = user.role
                db.session.commit()
        musician_choices.append((musician.id, user.get_display_name()))
    form.musician_id.choices = musician_choices
    
    if form.validate_on_submit():
        assignment = PracticeMusician(
            practice_id=practice_id,
            musician_id=form.musician_id.data,
            instrument=form.instrument.data
        )
        db.session.add(assignment)
        db.session.flush()  # Flush to get assignment ID
        
        # Get the musician object
        musician = Musician.query.get(form.musician_id.data)
        
        # Send SMS notification
        sms_sent = False
        sms_error = None
        if musician:
            try:
                result = send_practice_assignment_sms(practice, musician, is_new_assignment=True)
                # Handle old (success, error), new (success, error, sid), and latest (success, error, sid, status) formats
                if len(result) == 4:
                    success, error, message_sid, twilio_status = result
                elif len(result) == 3:
                    success, error, message_sid = result
                    twilio_status = None
                else:
                    success, error = result[0], result[1] if len(result) > 1 else None
                    message_sid = None
                    twilio_status = None
                
                # Get user info for logging
                user = musician.user if musician.user_id else None
                recipient_name = user.get_display_name() if user else musician.get_display_name()
                recipient_phone = format_phone_number(user.mobile_number) if user and user.mobile_number else None
                
                # Log SMS attempt
                try:
                    sms_log = SMSLog(
                        recipient_user_id=user.id if user else None,
                        recipient_phone=recipient_phone or 'Unknown',
                        recipient_name=recipient_name,
                        message_type='practice_assignment',
                        practice_id=practice.id,
                        musician_id=musician.id,
                        message_content=f"Practice assignment notification for {practice.date.strftime('%B %d, %Y') if practice.date else 'TBD'}",
                        status='success' if success else 'failed',
                        twilio_status=twilio_status,
                        error_message=error if not success else None,
                        twilio_message_sid=message_sid,
                        sent_by_user_id=current_user.id
                    )
                    db.session.add(sms_log)
                except Exception as log_error:
                    print(f"Warning: Could not log SMS: {log_error}")
                
                if success:
                    sms_sent = True
                    # Schedule SMS reminders (1 day before and 1 hour before)
                    if practice.date and practice.time:
                        try:
                            schedule_practice_sms_reminders(practice, musician)
                        except Exception as e:
                            print(f"Warning: Could not schedule SMS reminders: {e}")
                    
                    # Commit before redirecting
                    db.session.commit()
                    
                    # For admins/worship leaders, redirect to success page
                    if current_user.is_admin() or current_user.is_worship_leader():
                        return redirect(url_for('sms_success', practice_id=practice_id, musician_id=musician.id))
                else:
                    sms_error = error
                    flash(f'Team member added to practice. SMS notification failed: {error}', 'warning')
            except Exception as e:
                sms_error = str(e)
                flash(f'Team member added to practice. SMS notification error: {str(e)}', 'warning')
        
        db.session.commit()
        
        # Schedule SMS reminders (1 day before and 1 hour before) - fallback if SMS failed
        if musician and practice.date and practice.time and not sms_sent:
            try:
                schedule_practice_sms_reminders(practice, musician)
            except Exception as e:
                print(f"Warning: Could not schedule SMS reminders: {e}")
        
        # Expire the practice object's musicians relationship so it reloads on next access
        db.session.expire(practice, ['musicians'])
        
        if not musician or not any('SMS' in str(m) for m in getattr(flash, '_messages', [])):
            flash('Team member added to practice.', 'success')
    else:
        if form.errors:
            flash(f'Error adding team member: {form.errors}', 'danger')
        else:
            flash('Error adding team member. Please check the form.', 'danger')
    return redirect(url_for('practice_detail', id=practice_id))


@app.route('/sms-success')
@login_required
def sms_success():
    """Display SMS success page for admins/worship leaders"""
    if not (current_user.is_admin() or current_user.is_worship_leader()):
        flash('Access denied.', 'danger')
        return redirect(url_for('dashboard'))
    
    practice_id = request.args.get('practice_id', type=int)
    musician_id = request.args.get('musician_id', type=int)
    
    if not practice_id or not musician_id:
        flash('Invalid request.', 'danger')
        return redirect(url_for('practices'))
    
    practice = Practice.query.get_or_404(practice_id)
    musician = Musician.query.get_or_404(musician_id)
    user = musician.user if musician.user_id else None
    
    return render_template('sms_success.html', 
                         practice=practice, 
                         musician=musician,
                         user=user)


@app.route('/practices/<int:practice_id>/musicians/<int:assignment_id>/delete', methods=['POST'])
@worship_leader_required
def delete_practice_musician(practice_id, assignment_id):
    assignment = PracticeMusician.query.get_or_404(assignment_id)
    db.session.delete(assignment)
    db.session.commit()
    flash('Team member removed from practice.', 'success')
    return redirect(url_for('practice_detail', id=practice_id))


@app.route('/practices/<int:practice_id>/songs/add', methods=['POST'])
@app.route('/practices/<int:practice_id>/songs/add-multiple', methods=['POST'])
@login_required
def add_practice_songs(practice_id):
    """Add multiple songs to practice lineup at once"""
    practice = Practice.query.get_or_404(practice_id)
    
    # Get form data - handle both single and multiple songs
    song_names = request.form.getlist('song_names[]')
    keys = request.form.getlist('keys[]')
    speeds = request.form.getlist('speeds[]')
    orders = request.form.getlist('orders[]')
    prepared_by = request.form.get('prepared_by', type=int)
    
    # Fallback for single song (backward compatibility)
    if not song_names:
        song_name = request.form.get('song_name', '').strip()
        if song_name:
            song_names = [song_name]
            keys = [request.form.get('key', '').strip()]
            speeds = [request.form.get('speed', '')]
            orders = [request.form.get('order', type=int) or 0]
    
    if not song_names or not any(s.strip() for s in song_names):
        flash('Please enter at least one song name.', 'warning')
        return redirect(url_for('practice_detail', id=practice_id))
    
    # Get current max order
    max_order = db.session.query(db.func.max(PracticeSong.order)).filter_by(practice_id=practice_id).scalar() or 0
    
    added_count = 0
    skipped_count = 0
    
    for i, song_name in enumerate(song_names):
        song_name = song_name.strip()
        if not song_name:
            continue
        
        # Get corresponding values (with defaults)
        key = keys[i].strip() if i < len(keys) and keys[i] else None
        speed = speeds[i] if i < len(speeds) and speeds[i] else None
        order = int(orders[i]) if i < len(orders) and orders[i] else (max_order + i + 1)
        
        # Check if song already in lineup (case-insensitive)
        existing = PracticeSong.query.filter_by(practice_id=practice_id).filter(
            db.func.lower(PracticeSong.song_name) == song_name.lower()
        ).first()
        
        if existing:
            skipped_count += 1
            continue
        
        practice_song = PracticeSong(
            practice_id=practice_id,
            song_id=None,  # Custom songs only
            song_name=song_name,
            key=key if key else None,
            speed=speed if speed else None,
            prepared_by=prepared_by if prepared_by else None,
            order=order
        )
        db.session.add(practice_song)
        added_count += 1
    
    db.session.commit()
    
    if added_count > 0:
        flash(f'{added_count} song(s) added to lineup.', 'success')
    if skipped_count > 0:
        flash(f'{skipped_count} song(s) skipped (already in lineup).', 'warning')
    
    return redirect(url_for('practice_detail', id=practice_id))


@app.route('/practices/<int:practice_id>/songs/<int:practice_song_id>/remove', methods=['POST'])
@login_required
def remove_practice_song(practice_id, practice_song_id):
    """Remove a song from practice lineup"""
    practice_song = PracticeSong.query.get_or_404(practice_song_id)
    if practice_song.practice_id != practice_id:
        abort(404)
    db.session.delete(practice_song)
    db.session.commit()
    flash('Song removed from lineup.', 'success')
    return redirect(url_for('practice_detail', id=practice_id))


def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']




# Slide routes (Admin only)
@app.route('/slides')
@login_required
def slides():
    from sqlalchemy import func  # type: ignore
    search_query = request.args.get('search', '').strip()
    selected_artist = request.args.get('artist', '').strip()
    
    try:
        # Start with base query
        query = Slide.query
        
        # Apply title search filter
        if search_query:
            search_filter = f'%{search_query.lower()}%'
            query = query.filter(func.lower(Slide.title).like(search_filter))
        
        # Apply artist filter
        if selected_artist:
            artist_filter = f'%{selected_artist.lower()}%'
            query = query.filter(func.lower(Slide.artist).like(artist_filter))
        
        slides_list = query.order_by(Slide.title).all()
        
    except Exception as e:
        # If error due to missing column, run migration and retry
        if 'gender_key' in str(e):
            migrate_database()
            query = Slide.query
            
            if search_query:
                search_filter = f'%{search_query.lower()}%'
                query = query.filter(func.lower(Slide.title).like(search_filter))
            
            if selected_artist:
                artist_filter = f'%{selected_artist.lower()}%'
                query = query.filter(func.lower(Slide.artist).like(artist_filter))
            
            slides_list = query.order_by(Slide.title).all()
        else:
            raise
    
    return render_template('slides.html', 
                         slides=slides_list, 
                         search_query=search_query,
                         selected_artist=selected_artist)


def allowed_slide_file(filename):
    """Check if slide file extension is allowed"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_SLIDE_EXTENSIONS']


def detect_file_type_from_extension(filename):
    """Detect file type from file extension"""
    if not filename:
        return None
    
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    # Word documents
    if ext in ['doc', 'docx']:
        return 'word'
    # Excel/Spreadsheets
    elif ext in ['xls', 'xlsx']:
        return 'excel'
    # CSV files
    elif ext == 'csv':
        return 'csv'
    # Images
    elif ext in ['jpg', 'jpeg', 'png', 'gif']:
        return 'image'
    # PDF
    elif ext == 'pdf':
        return 'pdf'
    # Text files
    elif ext == 'txt':
        return 'txt'
    # PowerPoint
    elif ext in ['ppt', 'pptx']:
        return 'powerpoint'
    else:
        return None


@app.route('/slides/add', methods=['GET', 'POST'])
@admin_required
def add_slide():
    form = SlideForm()
    # Populate contributor dropdown with users
    users = User.query.order_by(User.username).all()
    form.artist.choices = [('', 'Select Contributor...')] + [(user.get_display_name(), user.get_display_name()) for user in users]
    if form.validate_on_submit():
        file_path = None
        
        # Handle file upload
        if form.slide_file.data:
            file = form.slide_file.data
            if file and hasattr(file, 'filename') and allowed_slide_file(file.filename):
                # Generate secure filename
                filename = secure_filename(file.filename)
                # Add timestamp to avoid conflicts
                name, ext = os.path.splitext(filename)
                filename = f"{name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{ext}"
                
                # Save file
                slides_dir = app.config['SLIDES_FOLDER']
                os.makedirs(slides_dir, exist_ok=True)
                file_path_full = os.path.join(slides_dir, filename)
                file.save(file_path_full)
                file_path = filename
                flash(f'File uploaded: {filename}', 'info')
            else:
                flash('Invalid file type. Please upload Word, Excel, CSV, PowerPoint, PDF, TXT, or Image files only.', 'danger')
                return render_template('slide_form.html', form=form, title='Add Slide')
        
        # Use manual file path if no file uploaded
        if not file_path and form.file_path.data:
            file_path = form.file_path.data.strip()
        
        # Clean title (replace underscores with spaces)
        cleaned_title = clean_slide_title(form.title.data)
        
        # Automatically detect file type from uploaded file
        file_type = None
        if file_path:
            file_type = detect_file_type_from_extension(file_path)
        
        slide = Slide(
            title=cleaned_title,
            artist=form.artist.data or None,
            description=form.description.data or None,
            language=None,  # Language field removed
            file_type=file_type,
            file_path=file_path,
            created_by=current_user.id
        )
        db.session.add(slide)
        db.session.commit()
        
        # Log activity
        log_activity(
            activity_type='job_aid_uploaded',
            actor_id=current_user.id,
            description=f"{current_user.get_display_name()} uploaded a new job aid: {cleaned_title}",
            slide_id=slide.id,
            metadata={'title': cleaned_title, 'file_type': file_type or 'unknown', 'contributor': form.artist.data or 'N/A'}
        )
        
        flash('Slide added successfully.', 'success')
        return redirect(url_for('slides'))
    return render_template('slide_form.html', form=form, title='Add Slide')


@app.route('/slides/<int:id>/edit', methods=['GET', 'POST'])
@admin_required
def edit_slide(id):
    slide = Slide.query.get_or_404(id)
    form = SlideForm(obj=slide)
    # Populate contributor dropdown with users
    users = User.query.order_by(User.username).all()
    user_display_names = [user.get_display_name() for user in users]
    form.artist.choices = [('', 'Select Contributor...')] + [(name, name) for name in user_display_names]
    # If existing artist value doesn't match any user, add it as an option
    if slide.artist and slide.artist not in user_display_names:
        form.artist.choices.append((slide.artist, slide.artist))
    # Pre-populate file_path with "Title Job Aid" format if it exists
    if slide.file_path and not form.file_path.data:
        # Get extension from existing file
        _, ext = os.path.splitext(slide.file_path)
        # Format as "Title Job Aid.ext"
        file_path_value = f"{slide.title.replace('_', ' ')} Job Aid{ext}"
        form.file_path.data = file_path_value
    if form.validate_on_submit():
        file_path = slide.file_path  # Keep existing file path by default
        
        # Handle file upload
        if form.slide_file.data:
            file = form.slide_file.data
            if file and hasattr(file, 'filename') and allowed_slide_file(file.filename):
                # Generate secure filename
                filename = secure_filename(file.filename)
                # Add timestamp to avoid conflicts
                name, ext = os.path.splitext(filename)
                filename = f"{name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{ext}"
                
                # Delete old file if it exists
                if slide.file_path:
                    old_file_path = os.path.join(app.config['SLIDES_FOLDER'], slide.file_path)
                    if os.path.exists(old_file_path):
                        try:
                            os.remove(old_file_path)
                        except:
                            pass
                
                # Save new file
                slides_dir = app.config['SLIDES_FOLDER']
                os.makedirs(slides_dir, exist_ok=True)
                file_path_full = os.path.join(slides_dir, filename)
                file.save(file_path_full)
                file_path = filename
                flash(f'File uploaded: {filename}', 'info')
            else:
                flash('Invalid file type. Please upload Word, Excel, CSV, PowerPoint, PDF, TXT, or Image files only.', 'danger')
                return render_template('slide_form.html', form=form, title='Edit Slide', slide=slide)
        
        # Use manual file path if no file uploaded and path is provided
        # But don't use it if it's in the "Title Job Aid" format (display only)
        if not form.slide_file.data and form.file_path.data:
            provided_path = form.file_path.data.strip()
            # Check if it's a display format (contains "Job Aid" and matches title)
            if "Job Aid" in provided_path:
                # Keep existing file_path, don't update with display format
                file_path = slide.file_path
            else:
                # It's a real file path, use it
                file_path = provided_path
        
        # Clean title (replace underscores with spaces)
        cleaned_title = clean_slide_title(form.title.data)
        
        # Automatically detect file type from uploaded file
        # If a new file was uploaded, detect from new file; otherwise keep existing file_type
        if form.slide_file.data and file_path:
            file_type = detect_file_type_from_extension(file_path)
        elif file_path:
            # File path was manually entered, detect from that
            file_type = detect_file_type_from_extension(file_path)
        else:
            # No file path, keep existing file_type
            file_type = slide.file_type
        
        slide.title = cleaned_title
        slide.artist = form.artist.data or None
        slide.description = form.description.data or None
        slide.file_type = file_type
        slide.file_path = file_path
        db.session.commit()
        flash('Slide updated successfully.', 'success')
        return redirect(url_for('slides'))
    return render_template('slide_form.html', form=form, title='Edit Slide', slide=slide)


@app.route('/slides/<int:id>/delete', methods=['POST'])
@admin_required
def delete_slide(id):
    slide = Slide.query.get_or_404(id)
    
    # Delete file if exists
    if slide.file_path:
        file_path = os.path.join(app.config['SLIDES_FOLDER'], slide.file_path)
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except:
                pass
    
    db.session.delete(slide)
    db.session.commit()
    flash('Slide deleted successfully.', 'success')
    return redirect(url_for('slides'))


def detect_language_from_title(title):
    """Detect language based on title content - defaults to Tagalog if not clearly English"""
    if not title:
        return 'tagalog'  # Default to Tagalog
    
    title_lower = title.lower().strip()
    
    # Common English worship-related words and phrases
    # If title contains these, it's likely English
    english_indicators = [
        'amazing', 'grace', 'great', 'good', 'love', 'lord', 'god', 'jesus', 'christ',
        'holy', 'spirit', 'praise', 'worship', 'bless', 'blessed', 'blessing',
        'glory', 'glorious', 'power', 'mighty', 'king', 'kingdom', 'heaven',
        'earth', 'forever', 'eternal', 'faith', 'hope', 'peace', 'joy', 'light',
        'darkness', 'victory', 'salvation', 'redeem', 'redeemer', 'savior', 'save',
        'cross', 'blood', 'sacrifice', 'resurrection', 'risen', 'ascend', 'come',
        'go', 'walk', 'run', 'stand', 'sit', 'sing', 'dance', 'shout', 'cry',
        'pray', 'prayer', 'thank', 'thanks', 'grateful', 'gratitude', 'mercy',
        'mercies', 'compassion', 'kindness', 'gentle', 'humble', 'meek', 'strong',
        'strength', 'weak', 'weakness', 'heal', 'healing', 'restore', 'restoration',
        'revive', 'revival', 'awaken', 'awakening', 'breakthrough', 'break', 'through',
        'freedom', 'free', 'liberty', 'deliver', 'deliverance', 'protect', 'protection',
        'shield', 'refuge', 'shelter', 'fortress', 'tower', 'rock', 'stone', 'cornerstone',
        'foundation', 'corner', 'stone', 'pillar', 'anchor', 'hope', 'trust', 'believe',
        'faith', 'faithful', 'faithfulness', 'true', 'truth', 'honest', 'honesty',
        'pure', 'purity', 'clean', 'cleanse', 'wash', 'white', 'snow', 'wool',
        'lamb', 'sheep', 'shepherd', 'pastor', 'flock', 'fold', 'green', 'pasture',
        'still', 'waters', 'valley', 'shadow', 'death', 'rod', 'staff', 'comfort',
        'table', 'presence', 'enemies', 'anoint', 'anointing', 'oil', 'cup', 'overflow',
        'goodness', 'mercy', 'follow', 'dwell', 'house', 'forever', 'eternity',
        'hillsong', 'bethel', 'elevation', 'passion', 'tomlin', 'redman', 'hughes',
        'stanfill', 'carnes', 'maher', 'baloche', 'townend', 'getty', 'wickham',
        'mullins', 'chapman', 'grant', 'crouch', 'gaither', 'smith', 'wesley',
        'newton', 'crosby', 'fanny', 'crosby', 'watts', 'isaac', 'wesley', 'charles'
    ]
    
    # Common Tagalog/Filipino worship-related words and phrases
    tagalog_indicators = [
        'araw', 'tahanan', 'panginoon', 'katapatan', 'binabago', 'baliw',
        'banal', 'dakilang', 'salamat', 'puri', 'awit', 'dasal', 'pag-ibig',
        'biyaya', 'kapangyarihan', 'kaluwalhatian', 'kabanalan', 'panalangin',
        'pagsamba', 'papuri', 'kagalakan', 'kapayapaan', 'kaligtasan',
        'katulad', 'tulad', 'walang', 'mayroon', 'ninyo', 'nila', 'namin',
        'iyan', 'iyon', 'dito', 'doon', 'kapag', 'dahil', 'kaya', 'opo',
        'sige', 'tama', 'mali', 'hindi', 'oo', 'hindi po', 'ng', 'sa', 'ang',
        'mga', 'na', 'ay', 'ko', 'mo', 'niya', 'ito', 'kung', 'pero', 'at', 'o'
    ]
    
    # First check for explicit Tagalog words
    for indicator in tagalog_indicators:
        if indicator in title_lower:
            return 'tagalog'
    
    # Then check for English words - if found, it's English
    for indicator in english_indicators:
        if indicator in title_lower:
            return 'english'
    
    # Default to Tagalog if not clearly English
    return 'tagalog'


def clean_slide_title(title):
    """Clean slide title by replacing underscores with spaces"""
    if not title:
        return title
    return title.replace('_', ' ').strip()


@app.route('/slides/bulk-change-language', methods=['POST'])
@admin_required
def bulk_change_language():
    """Bulk change language for multiple slides"""
    language = request.form.get('language')
    slide_ids_str = request.form.get('slide_ids', '')
    
    if not language:
        flash('Please select a language.', 'danger')
        return redirect(url_for('slides'))
    
    if not slide_ids_str:
        flash('Please select at least one slide.', 'warning')
        return redirect(url_for('slides'))
    
    try:
        slide_ids = [int(id.strip()) for id in slide_ids_str.split(',') if id.strip()]
        
        if not slide_ids:
            flash('No valid slides selected.', 'warning')
            return redirect(url_for('slides'))
        
        # Update language for selected slides
        updated_count = 0
        for slide_id in slide_ids:
            slide = Slide.query.get(slide_id)
            if slide:
                slide.language = language
                updated_count += 1
        
        db.session.commit()
        flash(f'Successfully updated language to {language.title()} for {updated_count} slide(s).', 'success')
    except Exception as e:
        db.session.rollback()
        flash(f'Error updating languages: {str(e)}', 'danger')
    
    return redirect(url_for('slides'))


@app.route('/slides/fix-titles', methods=['POST'])
@admin_required
def fix_slide_titles():
    """Fix all slide titles: replace underscores with spaces and auto-detect language"""
    slides = Slide.query.all()
    updated_count = 0
    title_updates = 0
    language_updates = 0
    
    for slide in slides:
        original_title = slide.title
        cleaned_title = clean_slide_title(original_title)
        title_changed = False
        
        # Update title if it had underscores
        if cleaned_title != original_title:
            slide.title = cleaned_title
            title_changed = True
            title_updates += 1
        
        # Always re-detect and set language based on cleaned title
        detected_language = detect_language_from_title(cleaned_title)
        if slide.language != detected_language:
            slide.language = detected_language
            language_updates += 1
        
        if title_changed or slide.language != detected_language:
            updated_count += 1
    
    if updated_count > 0:
        db.session.commit()
        messages = []
        if title_updates > 0:
            messages.append(f'{title_updates} title(s) cleaned')
        if language_updates > 0:
            messages.append(f'{language_updates} language(s) updated')
        flash(f'Successfully updated {updated_count} slide(s): {", ".join(messages)}.', 'success')
    else:
        flash('No slides needed updating.', 'info')
    
    return redirect(url_for('slides'))


# Event Announcement routes
@app.route('/notifications-page')
@login_required
def notifications_page():
    """Notifications page showing all notifications"""
    # Force fresh query to avoid caching issues
    notifications = Notification.query.filter_by(user_id=current_user.id).order_by(Notification.created_at.desc()).all()
    unread_count = Notification.query.filter_by(user_id=current_user.id, is_read=False).count()
    
    # Format notifications for display
    notifications_list = []
    for notif in notifications:
        actor_name = notif.actor.get_display_name() if notif.actor else 'Someone'
        notification_text = ''
        link = '#'
        icon = '🔔'
        
        if notif.notification_type == 'like':
            notification_text = f"{actor_name} liked your post"
            icon = '👍'
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'heart':
            notification_text = f"{actor_name} ❤️ your post"
            icon = '❤️'
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'share':
            notification_text = f"{actor_name} shared your post"
            icon = '🔄'
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'comment':
            notification_text = f"{actor_name} commented on your post"
            icon = '💬'
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'practice':
            notification_text = f"{actor_name} created a new practice schedule"
            icon = '📅'
            if notif.practice:
                link = url_for('practice_detail', id=notif.practice.id)
        elif notif.notification_type == 'leave_request':
            notification_text = f"{actor_name} filed a leave request for your approval"
            icon = '📋'
            link = '#'  # Will be handled by JavaScript to show popup
        elif notif.notification_type == 'leave_approved':
            notification_text = f"Your leave request has been approved by {actor_name}"
            icon = '✅'
            link = url_for('leave_requests')
        elif notif.notification_type == 'leave_rejected':
            # Include rejection reason if available
            if notif.leave_request and notif.leave_request.review_notes:
                notification_text = f"Your leave request has been rejected by {actor_name}. Reason: {notif.leave_request.review_notes}"
            else:
                notification_text = f"Your leave request has been rejected by {actor_name}"
            icon = '❌'
            link = url_for('leave_requests')
        
        notification_item = {
            'id': notif.id,
            'icon': icon,
            'text': notification_text,
            'link': link,
            'is_read': notif.is_read,
            'created_at': notif.created_at,
            'time_ago': _time_ago(notif.created_at),
            'type': notif.notification_type
        }
        
        # Add leave request ID if applicable
        if notif.leave_request_id:
            notification_item['leave_request_id'] = notif.leave_request_id
        
        notifications_list.append(notification_item)
    
    response = make_response(render_template('notifications.html', 
                         notifications=notifications_list, 
                         unread_count=unread_count))
    # Prevent caching to ensure fresh data
    response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    return response


@app.route('/announcements')
@admin_required
def announcements():
    """List all event announcements"""
    announcements_list = EventAnnouncement.query.order_by(EventAnnouncement.display_order, EventAnnouncement.created_at.desc()).all()
    return render_template('announcements.html', announcements=announcements_list)


@app.route('/announcements/add', methods=['GET', 'POST'])
@admin_required
def add_announcement():
    """Add a new event announcement"""
    form = EventAnnouncementForm()
    
    if form.validate_on_submit():
        # Handle file upload
        image_path = None
        if form.image.data:
            file = form.image.data
            if file.filename:
                filename = secure_filename(file.filename)
                # Add timestamp to avoid conflicts
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                name, ext = os.path.splitext(filename)
                filename = f"{name}_{timestamp}{ext}"
                
                announcements_dir = app.config['ANNOUNCEMENTS_FOLDER']
                os.makedirs(announcements_dir, exist_ok=True)
                file_path_full = os.path.join(announcements_dir, filename)
                file.save(file_path_full)
                image_path = f"announcements/{filename}"
                flash(f'Image uploaded: {filename}', 'info')
        
        # Parse display_order
        display_order = 0
        if form.display_order.data:
            try:
                display_order = int(form.display_order.data)
            except ValueError:
                display_order = 0
        
        # Parse is_active
        is_active = form.is_active.data == 'True'
        
        announcement = EventAnnouncement(
            title=form.title.data,
            caption=form.caption.data or None,
            image_path=image_path,
            is_active=is_active,
            display_order=display_order,
            created_by=current_user.id
        )
        db.session.add(announcement)
        db.session.commit()
        flash('Announcement added successfully.', 'success')
        return redirect(url_for('announcements'))
    
    return render_template('announcement_form.html', form=form, title='Add Event Announcement')


@app.route('/announcements/<int:id>/edit', methods=['GET', 'POST'])
@admin_required
def edit_announcement(id):
    """Edit an existing event announcement"""
    announcement = EventAnnouncement.query.get_or_404(id)
    form = EventAnnouncementForm(obj=announcement)
    
    # Set form defaults
    form.is_active.data = str(announcement.is_active)
    form.display_order.data = str(announcement.display_order) if announcement.display_order else '0'
    
    if form.validate_on_submit():
        # Handle file upload
        if form.image.data:
            file = form.image.data
            if file.filename:
                # Delete old file if exists
                if announcement.image_path:
                    old_file_path = os.path.join(app.config['ANNOUNCEMENTS_FOLDER'], os.path.basename(announcement.image_path))
                    if os.path.exists(old_file_path):
                        try:
                            os.remove(old_file_path)
                        except:
                            pass
                
                filename = secure_filename(file.filename)
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                name, ext = os.path.splitext(filename)
                filename = f"{name}_{timestamp}{ext}"
                
                announcements_dir = app.config['ANNOUNCEMENTS_FOLDER']
                os.makedirs(announcements_dir, exist_ok=True)
                file_path_full = os.path.join(announcements_dir, filename)
                file.save(file_path_full)
                announcement.image_path = f"announcements/{filename}"
                flash(f'Image uploaded: {filename}', 'info')
        
        announcement.title = form.title.data
        announcement.caption = form.caption.data or None
        
        # Parse display_order
        if form.display_order.data:
            try:
                announcement.display_order = int(form.display_order.data)
            except ValueError:
                announcement.display_order = 0
        
        # Parse is_active
        announcement.is_active = form.is_active.data == 'True'
        
        db.session.commit()
        flash('Announcement updated successfully.', 'success')
        return redirect(url_for('announcements'))
    
    return render_template('announcement_form.html', form=form, title='Edit Event Announcement', announcement=announcement)


@app.route('/announcements/<int:id>/toggle', methods=['POST'])
@admin_required
def toggle_announcement(id):
    """Toggle announcement active status"""
    announcement = EventAnnouncement.query.get_or_404(id)
    announcement.is_active = not announcement.is_active
    db.session.commit()
    
    status = 'activated' if announcement.is_active else 'deactivated'
    flash(f'Announcement {status} successfully.', 'success')
    return redirect(url_for('announcements'))


@app.route('/announcements/<int:id>/delete', methods=['POST'])
@admin_required
def delete_announcement(id):
    """Delete an event announcement"""
    announcement = EventAnnouncement.query.get_or_404(id)
    
    # Delete file if exists
    if announcement.image_path:
        file_path = os.path.join(app.config['ANNOUNCEMENTS_FOLDER'], os.path.basename(announcement.image_path))
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except:
                pass
    
    db.session.delete(announcement)
    db.session.commit()
    flash('Announcement deleted successfully.', 'success')
    return redirect(url_for('announcements'))



@app.route('/slides/<int:id>/view')
@login_required
def view_slide(id):
    """View file content - supports multiple file types"""
    slide = Slide.query.get_or_404(id)
    if not slide.file_path:
        flash('No Job Aid file available for this slide.', 'warning')
        return redirect(url_for('slides'))
    
    # Get filename - stored file_path should just be the filename
    filename = os.path.basename(slide.file_path)
    slides_dir = app.config['SLIDES_FOLDER']
    file_path = os.path.join(slides_dir, filename)
    
    if not os.path.exists(file_path):
        # Try alternative path construction
        alt_path = os.path.join(app.root_path, 'static', 'slides', filename)
        if os.path.exists(alt_path):
            file_path = alt_path
        else:
            flash(f'Job Aid file not found. Looking for: {filename} in {slides_dir}', 'danger')
            return redirect(url_for('slides'))
    
    # Get file extension
    _, ext = os.path.splitext(filename)
    ext = ext.lower()
    
    content = None
    content_type = None
    file_type = slide.file_type or detect_file_type_from_extension(filename)
    
    try:
        if file_type == 'word' or ext in ['.doc', '.docx']:
            # Read Word document - extract ALL content including tables, images, formatting
            try:
                from docx import Document  # type: ignore
                
                doc = Document(file_path)
                content_parts = []
                
                # Extract ALL paragraphs in order (Word docs don't have "pages" - just paragraphs)
                # This will get ALL content from the entire document
                for para in doc.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        content_parts.append(para_text)
                
                # Extract ALL tables (tables exist separately from paragraphs)
                for idx, table in enumerate(doc.tables):
                    content_parts.append(f"\n{'='*60}\nTABLE {idx + 1}\n{'='*60}")
                    
                    # Extract table with better formatting
                    for row in table.rows:
                        row_cells = []
                        for cell in row.cells:
                            # Get all text from cell, preserving content
                            cell_text_parts = []
                            for cell_para in cell.paragraphs:
                                if cell_para.text.strip():
                                    cell_text_parts.append(cell_para.text.strip())
                            cell_text = ' | '.join(cell_text_parts) if len(cell_text_parts) > 1 else (cell_text_parts[0] if cell_text_parts else '')
                            row_cells.append(cell_text)
                        
                        # Only add non-empty rows
                        if any(cell for cell in row_cells):
                            # Format as table row
                            table_row = ' | '.join(row_cells)
                            content_parts.append(table_row)
                    
                    content_parts.append(f"\n{'='*60}\n")
                
                # Extract headers from all sections
                headers = []
                try:
                    for section in doc.sections:
                        if section.header:
                            for para in section.header.paragraphs:
                                if para.text.strip():
                                    headers.append(para.text.strip())
                except:
                    pass
                
                # Extract footers from all sections
                footers = []
                try:
                    for section in doc.sections:
                        if section.footer:
                            for para in section.footer.paragraphs:
                                if para.text.strip():
                                    footers.append(para.text.strip())
                except:
                    pass
                
                # Combine all content with headers first, then body, then footers
                full_content = []
                if headers:
                    full_content.append("HEADERS:\n" + "\n".join(headers) + "\n" + "-"*60 + "\n")
                
                full_content.append("\n".join(content_parts))
                
                if footers:
                    full_content.append("\n" + "-"*60 + "\nFOOTERS:\n" + "\n".join(footers))
                
                content = "\n".join(full_content)
                content_type = 'text'
                    
            except ImportError:
                content = "Error: python-docx library not installed. Cannot read Word documents."
                content_type = 'error'
            except Exception as e:
                import traceback
                content = f"Error reading Word document: {str(e)}\n{traceback.format_exc()}"
                content_type = 'error'
        
        # If we have Word/Excel/PDF, provide file URL for embedded viewer
        if content_type == 'text' and file_type in ['word', 'excel', 'pdf']:
            file_url = url_for('static', filename=f'slides/{filename}')
            return render_template('view_slide.html', slide=slide, content=content, content_type=content_type, 
                                 file_type=file_type, file_url=file_url, embedded_view=True)
        
        elif file_type == 'excel' or ext in ['.xls', '.xlsx']:
            # Read Excel file - use openpyxl (more reliable, no compilation needed)
            try:
                import openpyxl  # type: ignore
                wb = openpyxl.load_workbook(file_path)
                content_parts = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    content_parts.append(f"=== Sheet: {sheet_name} ===\n")
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [str(cell) if cell is not None else '' for cell in row]
                        if any(cell for cell in row_data):  # Skip empty rows
                            content_parts.append('\t'.join(row_data))
                    content_parts.append("\n")
                content = '\n'.join(content_parts)
                content_type = 'text'
            except ImportError:
                content = "Error: openpyxl library not installed. Cannot read Excel files."
                content_type = 'error'
            except Exception as e:
                content = f"Error reading Excel file: {str(e)}"
                content_type = 'error'
        
        elif file_type == 'csv' or ext == '.csv':
            # Read CSV file
            try:
                import csv
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = []
                    for row in reader:
                        rows.append('\t'.join(row))
                    content = '\n'.join(rows)
                content_type = 'text'
            except Exception as e:
                content = f"Error reading CSV file: {str(e)}"
                content_type = 'error'
        
        elif file_type == 'pdf' or ext == '.pdf':
            # Read PDF file
            try:
                import PyPDF2  # type: ignore
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    text_parts = []
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        text = page.extract_text()
                        if text.strip():
                            text_parts.append(f"--- Page {page_num} ---\n{text}")
                    content = '\n\n'.join(text_parts)
                content_type = 'text'
            except ImportError:
                content = "Error: PyPDF2 library not installed. Cannot read PDF files."
                content_type = 'error'
            except Exception as e:
                content = f"Error reading PDF file: {str(e)}"
                content_type = 'error'
        
        elif file_type == 'txt' or ext == '.txt':
            # Read text file
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                content_type = 'text'
            except Exception as e:
                content = f"Error reading text file: {str(e)}"
                content_type = 'error'
        
        elif file_type == 'image' or ext in ['.jpg', '.jpeg', '.png', '.gif']:
            # Display image
            content = url_for('static', filename=f'slides/{filename}')
            content_type = 'image'
        
        elif file_type == 'powerpoint' or ext in ['.ppt', '.pptx']:
            # Read PowerPoint
            try:
                from pptx import Presentation  # type: ignore
                prs = Presentation(file_path)
                slide_texts = []
                for idx, slide_obj in enumerate(prs.slides, 1):
                    for shape in slide_obj.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                content = '\n\n'.join(slide_texts)
                content_type = 'text'
            except Exception as e:
                content = f"Error reading PowerPoint file: {str(e)}"
                content_type = 'error'
        
        else:
            content = f"File type '{ext}' is not supported for viewing."
            content_type = 'error'
    
    except Exception as e:
        flash(f'Error reading Job Aid file: {str(e)}', 'danger')
        return redirect(url_for('slides'))
    
    # For Word/Excel/PDF/PowerPoint files, also provide file URL for embedded viewing option
    file_url = None
    embedded_view = False
    if file_type in ['word', 'excel', 'pdf', 'powerpoint'] and content_type == 'text':
        # Get full URL for the file
        file_url = url_for('static', filename=f'slides/{filename}', _external=True)
        embedded_view = True
    
    return render_template('view_slide.html', slide=slide, content=content, content_type=content_type, 
                         file_type=file_type, file_url=file_url, embedded_view=embedded_view)


@app.route('/slides/<int:id>/download')
@login_required
def download_slide(id):
    slide = Slide.query.get_or_404(id)
    if not slide.file_path:
        flash('No Job Aid file available for this slide.', 'warning')
        return redirect(url_for('slides'))
    
    # Sanitize file path to prevent directory traversal
    original_filename = secure_filename(os.path.basename(slide.file_path))
    slides_dir = app.config['SLIDES_FOLDER']
    file_path = os.path.join(slides_dir, original_filename)
    
    if not os.path.exists(file_path):
        flash('Job Aid file not found.', 'danger')
        return redirect(url_for('slides'))
    
    # Generate custom filename: Task Name + " Job Aid" + extension
    task_name = slide.title.replace('_', ' ')
    # Get file extension from original filename
    _, file_ext = os.path.splitext(original_filename)
    # Create custom download filename
    custom_filename = f"{secure_filename(task_name)} Job Aid{file_ext}"
    
    return send_from_directory(slides_dir, original_filename, as_attachment=True, download_name=custom_filename)


# User Management routes (Admin only)
@app.route('/permissions', methods=['GET', 'POST'])
@admin_required
def permissions():
    """Manage user permissions - Admin only"""
    form = PermissionForm()
    
    # Populate user choices
    users = User.query.filter(User.role != 'admin').order_by(User.username).all()
    form.user_id.choices = [(u.id, f"{u.get_display_name()} ({u.username})") for u in users]
    
    if form.validate_on_submit():
        user = User.query.get(form.user_id.data)
        if not user:
            flash('User not found.', 'danger')
            return redirect(url_for('permissions'))
        
        # Define permission types
        permission_types = {
            'edit_slides': form.edit_slides.data,
            'edit_announcements': form.edit_announcements.data
        }
        
        # Update permissions
        for perm_type, granted in permission_types.items():
            existing = UserPermission.query.filter_by(
                user_id=user.id,
                permission_type=perm_type
            ).first()
            
            if granted:
                # Grant permission
                if not existing:
                    permission = UserPermission(
                        user_id=user.id,
                        permission_type=perm_type,
                        granted_by=current_user.id
                    )
                    db.session.add(permission)
            else:
                # Revoke permission
                if existing:
                    db.session.delete(existing)
        
        db.session.commit()
        flash(f'Permissions updated for {user.get_display_name()}.', 'success')
        return redirect(url_for('permissions'))
    
    # Get all users with their permissions
    all_users = User.query.filter(User.role != 'admin').order_by(User.username).all()
    users_with_permissions = []
    for user in all_users:
        permissions_dict = {p.permission_type: True for p in user.permissions}
        users_with_permissions.append({
            'user': user,
            'permissions': permissions_dict
        })
    
    return render_template('permissions.html', form=form, users_with_permissions=users_with_permissions)


@app.route('/tasks', methods=['GET', 'POST'])
@login_required
def tasks():
    """Tasks page with My Task Today and EOD Tasks"""
    today = date.today()
    
    # Get today's tasks (pending and completed)
    today_tasks = Task.query.filter_by(
        user_id=current_user.id,
        task_date=today
    ).order_by(Task.priority.asc(), Task.created_at.asc()).all()
    
    # Separate pending and completed tasks
    pending_tasks = [t for t in today_tasks if not t.is_completed]
    completed_tasks_today = [t for t in today_tasks if t.is_completed]
    
    # Get completed tasks for EOD summary (today's completed tasks)
    eod_tasks = completed_tasks_today
    
    # Get current date in Manila timezone
    import pytz  # type: ignore
    manila_tz = pytz.timezone('Asia/Manila')
    current_date_manila = datetime.now(manila_tz).strftime('%B %d, %Y')
    
    # Get user's saved task options
    task_options = TaskOption.query.filter_by(user_id=current_user.id).order_by(TaskOption.created_at.desc()).all()
    
    return render_template('tasks.html', 
                         pending_tasks=pending_tasks,
                         eod_tasks=eod_tasks,
                         current_date_manila=current_date_manila,
                         today=today,
                         task_options=task_options)


@app.route('/tools', methods=['GET'])
@login_required
def tools():
    """Tools page showing all tools"""
    try:
        tools_list = Tool.query.order_by(Tool.created_at.desc()).all()
    except Exception as e:
        # If error due to missing column, run migration and retry
        if 'developer_name' in str(e) or 'no such column' in str(e).lower():
            # Run migration using direct SQLite connection to ensure it works
            import sqlite3
            
            db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'instance', 'database.db')
            if os.path.exists(db_path):
                try:
                    sqlite_conn = sqlite3.connect(db_path)
                    cursor = sqlite_conn.cursor()
                    cursor.execute("PRAGMA table_info(tool)")
                    columns = [row[1] for row in cursor.fetchall()]
                    
                    if 'developer_name' not in columns:
                        cursor.execute("ALTER TABLE tool ADD COLUMN developer_name VARCHAR(200)")
                        sqlite_conn.commit()
                    sqlite_conn.close()
                except Exception as migration_error:
                    print(f"Migration error: {migration_error}")
            
            # Also run the standard migration
            migrate_database()
            
            # Retry the query
            tools_list = Tool.query.order_by(Tool.created_at.desc()).all()
        else:
            raise
    return render_template('tools.html', tools=tools_list)


@app.route('/tools/add', methods=['GET', 'POST'])
@login_required
def add_tool():
    """Add a new tool"""
    # Only allow admin or team leader to add tools
    if not current_user.is_admin() and not current_user.is_team_leader():
        flash('You do not have permission to add tools.', 'danger')
        return redirect(url_for('tools'))
    
    form = ToolForm()
    if form.validate_on_submit():
        # Handle screenshot upload
        screenshot_path = None
        if form.screenshot.data:
            file = form.screenshot.data
            if file and hasattr(file, 'filename') and file.filename:
                # Check if it's an allowed image type
                filename = file.filename.lower()
                if any(filename.endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif']):
                    # Generate secure filename
                    filename = secure_filename(f"tool_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                    
                    # Save file
                    tools_dir = app.config['TOOLS_FOLDER']
                    os.makedirs(tools_dir, exist_ok=True)
                    file_path_full = os.path.join(tools_dir, filename)
                    file.save(file_path_full)
                    screenshot_path = f"tools/{filename}"
                else:
                    flash('Invalid file type. Please upload JPG, PNG, or GIF images only.', 'danger')
                    return render_template('tool_form.html', form=form, title='Add Tool')
        
        tool = Tool(
            name=form.name.data,
            link=form.link.data,
            description=form.description.data or None,
            screenshot=screenshot_path,
            developer_name=form.developer_name.data or None,
            created_by=current_user.id
        )
        db.session.add(tool)
        db.session.commit()
        flash('Tool added successfully.', 'success')
        return redirect(url_for('tools'))
    return render_template('tool_form.html', form=form, title='Add Tool')


@app.route('/tools/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit_tool(id):
    """Edit an existing tool"""
    try:
        tool = Tool.query.get_or_404(id)
    except Exception as e:
        # If error due to missing column, run migration and retry
        if 'developer_name' in str(e) or 'no such column' in str(e).lower():
            # Run migration using direct SQLite connection to ensure it works
            import sqlite3
            
            db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'instance', 'database.db')
            if os.path.exists(db_path):
                try:
                    sqlite_conn = sqlite3.connect(db_path)
                    cursor = sqlite_conn.cursor()
                    cursor.execute("PRAGMA table_info(tool)")
                    columns = [row[1] for row in cursor.fetchall()]
                    
                    if 'developer_name' not in columns:
                        cursor.execute("ALTER TABLE tool ADD COLUMN developer_name VARCHAR(200)")
                        sqlite_conn.commit()
                    sqlite_conn.close()
                except Exception as migration_error:
                    print(f"Migration error: {migration_error}")
            
            # Also run the standard migration
            migrate_database()
            
            # Retry the query
            tool = Tool.query.get_or_404(id)
        else:
            raise
    
    # Only allow admin or team leader to edit
    if not current_user.is_admin() and not current_user.is_team_leader():
        flash('You do not have permission to edit this tool.', 'danger')
        return redirect(url_for('tools'))
    
    form = ToolForm(obj=tool)
    
    if form.validate_on_submit():
        tool.name = form.name.data
        tool.link = form.link.data
        tool.description = form.description.data or None
        tool.developer_name = form.developer_name.data or None
        
        # Handle screenshot upload
        if form.screenshot.data:
            file = form.screenshot.data
            if file and hasattr(file, 'filename') and file.filename:
                # Check if it's an allowed image type
                filename = file.filename.lower()
                if any(filename.endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif']):
                    # Delete old screenshot if exists
                    if tool.screenshot:
                        old_file_path = os.path.join(app.config['TOOLS_FOLDER'], os.path.basename(tool.screenshot))
                        if os.path.exists(old_file_path):
                            try:
                                os.remove(old_file_path)
                            except Exception as e:
                                print(f"Warning: Could not delete old screenshot: {e}")
                    
                    # Generate secure filename
                    filename = secure_filename(f"tool_{tool.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                    
                    # Save file
                    tools_dir = app.config['TOOLS_FOLDER']
                    os.makedirs(tools_dir, exist_ok=True)
                    file_path_full = os.path.join(tools_dir, filename)
                    file.save(file_path_full)
                    tool.screenshot = f"tools/{filename}"
                else:
                    flash('Invalid file type. Please upload JPG, PNG, or GIF images only.', 'danger')
                    return render_template('tool_form.html', form=form, title='Edit Tool', tool=tool)
        
        db.session.commit()
        flash('Tool updated successfully.', 'success')
        return redirect(url_for('tools'))
    return render_template('tool_form.html', form=form, title='Edit Tool', tool=tool)


@app.route('/tools/<int:id>/delete', methods=['POST'])
@login_required
def delete_tool(id):
    """Delete a tool"""
    try:
        tool = Tool.query.get_or_404(id)
    except Exception as e:
        # If error due to missing column, run migration and retry
        if 'developer_name' in str(e) or 'no such column' in str(e).lower():
            # Run migration using direct SQLite connection to ensure it works
            import sqlite3
            
            db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'instance', 'database.db')
            if os.path.exists(db_path):
                try:
                    sqlite_conn = sqlite3.connect(db_path)
                    cursor = sqlite_conn.cursor()
                    cursor.execute("PRAGMA table_info(tool)")
                    columns = [row[1] for row in cursor.fetchall()]
                    
                    if 'developer_name' not in columns:
                        cursor.execute("ALTER TABLE tool ADD COLUMN developer_name VARCHAR(200)")
                        sqlite_conn.commit()
                    sqlite_conn.close()
                except Exception as migration_error:
                    print(f"Migration error: {migration_error}")
            
            # Also run the standard migration
            migrate_database()
            
            # Retry the query
            tool = Tool.query.get_or_404(id)
        else:
            raise
    
    # Only allow admin or team leader to delete
    if not current_user.is_admin() and not current_user.is_team_leader():
        flash('You do not have permission to delete this tool.', 'danger')
        return redirect(url_for('tools'))
    
    # Delete screenshot if exists
    if tool.screenshot:
        screenshot_path = os.path.join(app.config['TOOLS_FOLDER'], os.path.basename(tool.screenshot))
        if os.path.exists(screenshot_path):
            try:
                os.remove(screenshot_path)
            except Exception as e:
                print(f"Warning: Could not delete screenshot: {e}")
    
    db.session.delete(tool)
    db.session.commit()
    flash('Tool deleted successfully.', 'success')
    return redirect(url_for('tools'))


@app.route('/tasks/add', methods=['POST'])
@csrf.exempt
@login_required
def add_task():
    """Add a new task"""
    try:
        data = request.get_json()
        task_text = data.get('task', '').strip()
        priority = int(data.get('priority', 1))
        task_date_str = data.get('date', date.today().isoformat())
        
        if not task_text:
            return jsonify({'success': False, 'message': 'Task description is required.'}), 400
        
        # Parse date
        try:
            task_date = datetime.strptime(task_date_str, '%Y-%m-%d').date()
        except ValueError:
            task_date = date.today()
        
        # Validate priority
        if priority not in [1, 2, 3]:
            priority = 1
        
        task = Task(
            user_id=current_user.id,
            task=task_text,
            priority=priority,
            task_date=task_date,
            is_completed=False
        )
        db.session.add(task)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': 'Task added successfully.',
            'task': {
                'id': task.id,
                'task': task.task,
                'priority': task.priority,
                'is_completed': task.is_completed
            }
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error adding task: {str(e)}'}), 500


@app.route('/tasks/<int:id>/complete', methods=['POST'])
@csrf.exempt
@login_required
def complete_task(id):
    """Mark a task as completed"""
    task = Task.query.get_or_404(id)
    
    if task.user_id != current_user.id:
        return jsonify({'success': False, 'message': 'You can only complete your own tasks.'}), 403
    
    task.is_completed = True
    task.completed_at = datetime.utcnow()
    db.session.commit()
    
    return jsonify({'success': True, 'message': 'Task marked as completed.'})


@app.route('/tasks/<int:id>/uncomplete', methods=['POST'])
@csrf.exempt
@login_required
def uncomplete_task(id):
    """Mark a task as not completed"""
    task = Task.query.get_or_404(id)
    
    if task.user_id != current_user.id:
        return jsonify({'success': False, 'message': 'You can only modify your own tasks.'}), 403
    
    task.is_completed = False
    task.completed_at = None
    db.session.commit()
    
    return jsonify({'success': True, 'message': 'Task marked as pending.'})


@app.route('/tasks/<int:id>/delete', methods=['POST'])
@csrf.exempt
@login_required
def delete_task(id):
    """Delete a task"""
    task = Task.query.get_or_404(id)
    
    if task.user_id != current_user.id:
        return jsonify({'success': False, 'message': 'You can only delete your own tasks.'}), 403
    
    db.session.delete(task)
    db.session.commit()
    
    return jsonify({'success': True, 'message': 'Task deleted successfully.'})


# Task Options routes
@app.route('/task-options/add', methods=['POST'])
@csrf.exempt
@login_required
def add_task_option():
    """Add a new task option/template"""
    try:
        data = request.get_json()
        task_text = data.get('task_text', '').strip()
        
        if not task_text:
            return jsonify({'success': False, 'message': 'Task description is required.'}), 400
        
        # Check if task option already exists for this user
        existing = TaskOption.query.filter_by(
            user_id=current_user.id,
            task_text=task_text
        ).first()
        
        if existing:
            return jsonify({'success': False, 'message': 'This task option already exists.'}), 400
        
        task_option = TaskOption(
            user_id=current_user.id,
            task_text=task_text,
            priority=2  # Default priority (medium)
        )
        db.session.add(task_option)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': 'Task option added successfully.',
            'task_option': {
                'id': task_option.id,
                'task_text': task_option.task_text
            }
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error adding task option: {str(e)}'}), 500


@app.route('/task-options/<int:id>/delete', methods=['POST'])
@csrf.exempt
@login_required
def delete_task_option(id):
    """Delete a task option"""
    task_option = TaskOption.query.get_or_404(id)
    
    if task_option.user_id != current_user.id:
        return jsonify({'success': False, 'message': 'You can only delete your own task options.'}), 403
    
    db.session.delete(task_option)
    db.session.commit()
    
    return jsonify({'success': True, 'message': 'Task option deleted successfully.'})


@app.route('/task-options/<int:id>/use', methods=['POST'])
@csrf.exempt
@login_required
def use_task_option(id):
    """Use a task option to create a new task"""
    task_option = TaskOption.query.get_or_404(id)
    
    if task_option.user_id != current_user.id:
        return jsonify({'success': False, 'message': 'You can only use your own task options.'}), 403
    
    try:
        data = request.get_json()
        task_date_str = data.get('date', date.today().isoformat())
        
        # Parse date
        try:
            task_date = datetime.strptime(task_date_str, '%Y-%m-%d').date()
        except ValueError:
            task_date = date.today()
        
        # Check if this task already exists for the user on this date
        existing_task = Task.query.filter_by(
            user_id=current_user.id,
            task=task_option.task_text,
            task_date=task_date
        ).first()
        
        if existing_task:
            return jsonify({
                'success': False,
                'message': 'This task already exists for today. Please use a different task option or add it manually.'
            }), 400
        
        task = Task(
            user_id=current_user.id,
            task=task_option.task_text,
            priority=2,  # Default to medium priority when using task option
            task_date=task_date,
            is_completed=False
        )
        db.session.add(task)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': 'Task added from option successfully.',
            'task': {
                'id': task.id,
                'task': task.task,
                'priority': task.priority,
                'is_completed': task.is_completed
            }
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error creating task: {str(e)}'}), 500


def format_task_as_sentence(task_text):
    """Convert a task name into a full sentence format.
    For example: 'Box contents' -> 'Processed Box Contents'
    """
    if not task_text:
        return ""
    
    # Common action verbs to prepend if task doesn't start with one
    action_verbs = ['processed', 'completed', 'reviewed', 'organized', 'updated', 'verified', 'checked']
    
    # Check if task already starts with a verb (common patterns)
    task_lower = task_text.lower().strip()
    starts_with_verb = any(task_lower.startswith(verb) for verb in action_verbs)
    
    # If it doesn't start with a verb, add "Processed"
    if not starts_with_verb:
        # Capitalize first letter of each word
        words = task_text.strip().split()
        capitalized_words = [word.capitalize() for word in words]
        formatted_task = " ".join(capitalized_words)
        return f"Processed {formatted_task}"
    else:
        # Already has a verb, just capitalize properly
        words = task_text.strip().split()
        capitalized_words = [word.capitalize() for word in words]
        return " ".join(capitalized_words)


@app.route('/tasks/eod/download', methods=['GET'])
@login_required
def download_eod_tasks():
    """Download completed tasks for today as a .txt file"""
    today = date.today()
    
    # Get today's completed tasks for current user
    completed_tasks = Task.query.filter_by(
        user_id=current_user.id,
        task_date=today,
        is_completed=True
    ).order_by(Task.completed_at.asc()).all()
    
    if not completed_tasks:
        flash('No completed tasks to download.', 'warning')
        return redirect(url_for('tasks'))
    
    # Get current date in Manila timezone for filename
    import pytz  # type: ignore
    manila_tz = pytz.timezone('Asia/Manila')
    current_date_manila = datetime.now(manila_tz).strftime('%Y%m%d')
    
    # Create text content with tasks as they are
    content_lines = []
    content_lines.append(f"EOD (End of Day) Tasks - {current_user.get_display_name()}")
    content_lines.append(f"Date: {datetime.now(manila_tz).strftime('%B %d, %Y')}")
    content_lines.append("")
    content_lines.append("Completed Tasks:")
    content_lines.append("")
    
    for task in completed_tasks:
        content_lines.append(task.task)
    
    content = "\n".join(content_lines)
    
    # Create response with text file
    filename = f"EOD_Tasks_{current_user.get_display_name().replace(' ', '_')}_{current_date_manila}.txt"
    
    response = Response(
        content,
        mimetype='text/plain',
        headers={
            'Content-Disposition': f'attachment; filename="{filename}"'
        }
    )
    
    return response


@app.route('/journal', methods=['GET', 'POST'])
@login_required
def journal():
    """Journal page with mood board only (prayers and devotions removed)"""
    form = JournalForm()
    
    # Get filter for entry type (only mood_board now)
    filter_type = request.args.get('type', 'mood_board')
    if filter_type not in ['mood_board', 'all']:
        filter_type = 'mood_board'
    
    # Build query
    query = Journal.query.filter_by(user_id=current_user.id)
    if filter_type != 'all':
        query = query.filter_by(entry_type='mood_board')
    
    # Get entries ordered by date (newest first)
    entries = query.order_by(Journal.date.desc(), Journal.created_at.desc()).all()
    
    # Group entries by type for display (only mood_board now)
    entries_by_type = {
        'mood_board': [e for e in entries if e.entry_type == 'mood_board']
    }
    
    if form.validate_on_submit():
        # Only allow mood_board entries
        if form.entry_type.data != 'mood_board':
            form.entry_type.data = 'mood_board'
        
        # Handle image upload for mood board
        image_path = None
        if form.image.data and form.entry_type.data == 'mood_board':
            file = form.image.data
            if hasattr(file, 'filename') and file.filename:
                filename = secure_filename(f"journal_{current_user.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['JOURNALS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['JOURNALS_FOLDER'], filename)
                file.save(file_path)
                image_path = f"journals/{filename}"
        
        # Parse date
        try:
            entry_date = datetime.strptime(form.date.data, '%Y-%m-%d').date()
        except ValueError:
            entry_date = date.today()
        
        # Get selected emojis for mood board
        mood_emojis = None
        if form.entry_type.data == 'mood_board' and form.mood_emojis.data:
            mood_emojis = form.mood_emojis.data
        
        # Create journal entry
        journal_entry = Journal(
            user_id=current_user.id,
            entry_type='mood_board',
            title=form.title.data or None,
            content=form.content.data or None,
            image_path=image_path,
            mood_emojis=mood_emojis,
            date=entry_date
        )
        db.session.add(journal_entry)
        db.session.commit()
        flash('Journal entry saved successfully!', 'success')
        return redirect(url_for('journal', type='mood_board'))
    
    # Set default date to today
    if not form.date.data:
        form.date.data = date.today().strftime('%Y-%m-%d')
    
    # Get current date in Manila timezone
    import pytz  # type: ignore
    manila_tz = pytz.timezone('Asia/Manila')
    current_date_manila = datetime.now(manila_tz).strftime('%B %d, %Y')
    
    return render_template('journal.html', form=form, entries=entries, entries_by_type=entries_by_type, filter_type=filter_type, current_date_manila=current_date_manila)


@app.route('/journal/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit_journal(id):
    """Edit a journal entry"""
    journal_entry = Journal.query.get_or_404(id)
    
    # Check if user owns this entry
    if journal_entry.user_id != current_user.id:
        flash('You do not have permission to edit this entry.', 'danger')
        return redirect(url_for('journal'))
    
    form = JournalForm(obj=journal_entry)
    form.date.data = journal_entry.date.strftime('%Y-%m-%d') if journal_entry.date else date.today().strftime('%Y-%m-%d')
    if journal_entry.mood_emojis:
        form.mood_emojis.data = journal_entry.mood_emojis
    if journal_entry.application:
        form.application.data = journal_entry.application
    if journal_entry.prayer_text:
        form.prayer_text.data = journal_entry.prayer_text
    
    if form.validate_on_submit():
        # Handle image upload
        if form.image.data and form.entry_type.data == 'mood_board':
            file = form.image.data
            if hasattr(file, 'filename') and file.filename:
                # Delete old image if exists
                if journal_entry.image_path:
                    old_path = os.path.join(app.root_path, 'static', journal_entry.image_path)
                    if os.path.exists(old_path):
                        os.remove(old_path)
                
                filename = secure_filename(f"journal_{current_user.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}")
                os.makedirs(app.config['JOURNALS_FOLDER'], exist_ok=True)
                file_path = os.path.join(app.config['JOURNALS_FOLDER'], filename)
                file.save(file_path)
                journal_entry.image_path = f"journals/{filename}"
        
        # Parse date
        try:
            journal_entry.date = datetime.strptime(form.date.data, '%Y-%m-%d').date()
        except ValueError:
            journal_entry.date = date.today()
        
        journal_entry.entry_type = form.entry_type.data
        journal_entry.title = form.title.data or None
        journal_entry.content = form.content.data or None
        
        # Update devotion-specific fields
        if form.entry_type.data == 'devotion':
            journal_entry.application = form.application.data or None
            journal_entry.prayer_text = form.prayer_text.data or None
        
        journal_entry.updated_at = datetime.utcnow()
        
        db.session.commit()
        flash('Journal entry updated successfully!', 'success')
        return redirect(url_for('journal', type=journal_entry.entry_type))
    
    return render_template('journal_form.html', form=form, journal_entry=journal_entry, title='Edit Journal Entry')


@app.route('/journal/<int:id>/delete', methods=['POST'])
@login_required
def delete_journal(id):
    """Delete a journal entry"""
    journal_entry = Journal.query.get_or_404(id)
    
    # Check if user owns this entry
    if journal_entry.user_id != current_user.id:
        flash('You do not have permission to delete this entry.', 'danger')
        return redirect(url_for('journal'))
    
    # Delete image if exists
    if journal_entry.image_path:
        image_path = os.path.join(app.root_path, 'static', journal_entry.image_path)
        if os.path.exists(image_path):
            os.remove(image_path)
    
    db.session.delete(journal_entry)
    db.session.commit()
    flash('Journal entry deleted successfully.', 'success')
    return redirect(url_for('journal'))


@app.route('/users')
@admin_required
def users():
    """List all users"""
    users_list = User.query.order_by(User.username).all()
    return render_template('users.html', users=users_list)


@app.route('/users/add', methods=['GET', 'POST'])
@admin_required
def add_user():
    """Add a new user"""
    form = UserForm()
    if form.validate_on_submit():
        # Check if username already exists
        existing_user = User.query.filter_by(username=form.username.data).first()
        if existing_user:
            flash('Username already exists. Please choose a different username.', 'danger')
            return render_template('user_form.html', form=form, title='Add User')
        
        # Check if email already exists
        existing_email = User.query.filter_by(email=form.email.data).first()
        if existing_email:
            flash('Email already exists. Please use a different email.', 'danger')
            return render_template('user_form.html', form=form, title='Add User')
        
        user = User(
            username=form.username.data,
            email=form.email.data,
            nickname=form.nickname.data.strip() if form.nickname.data else None,
            mobile_number=form.mobile_number.data.strip() if form.mobile_number.data else None,
            role=form.role.data
        )
        
        # Set password if provided, otherwise set default password
        if form.password.data and form.password.data.strip():
            user.set_password(form.password.data.strip())
        else:
            # Set default password if not provided (required for database NOT NULL constraint)
            user.set_password('password123')
            flash('Default password "password123" set. User should change it on first login.', 'warning')
        
        db.session.add(user)
        db.session.flush()  # Flush to get user.id
        
        # Auto-create musician profile with synced role
        musician = Musician(
            name=user.get_display_name(),
            user_id=user.id,
            instruments=user.role if user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader'] else None
        )
        db.session.add(musician)
        db.session.commit()
        flash('User created successfully.', 'success')
        return redirect(url_for('users'))
    
    return render_template('user_form.html', form=form, title='Add User')


@app.route('/users/<int:id>/edit', methods=['GET', 'POST'])
@admin_required
def edit_user(id):
    """Edit an existing user"""
    user = User.query.get_or_404(id)
    form = UserForm(obj=user)
    
    # Don't populate password field
    form.password.data = None
    
    if form.validate_on_submit():
        # Check if username is being changed and if new username exists
        if form.username.data != user.username:
            existing_user = User.query.filter_by(username=form.username.data).first()
            if existing_user:
                flash('Username already exists. Please choose a different username.', 'danger')
                return render_template('user_form.html', form=form, title='Edit User', user=user)
        
        # Check if email is being changed and if new email exists
        if form.email.data != user.email:
            existing_email = User.query.filter_by(email=form.email.data).first()
            if existing_email:
                flash('Email already exists. Please use a different email.', 'danger')
                return render_template('user_form.html', form=form, title='Edit User', user=user)
        
        user.username = form.username.data
        user.email = form.email.data
        user.nickname = form.nickname.data.strip() if form.nickname.data else None
        user.mobile_number = form.mobile_number.data.strip() if form.mobile_number.data else None
        old_role = user.role
        user.role = form.role.data
        
        # Update password if provided (optional - passwords disabled)
        if form.password.data and form.password.data.strip():
            user.set_password(form.password.data.strip())
            flash('Password updated.', 'info')
        
        # Sync role to musician profile if it changed
        if old_role != user.role and user.musician:
            if user.role in ['case_manager', 'shipment_coordinator', 'data_analyst', 'team_leader']:
                user.musician.instruments = user.role
            else:
                user.musician.instruments = None
        
        db.session.commit()
        flash('User updated successfully.', 'success')
        return redirect(url_for('users'))
    
    return render_template('user_form.html', form=form, title='Edit User', user=user)


@app.route('/users/<int:id>/reset-password', methods=['POST'])
@admin_required
def reset_user_password(id):
    """Reset a user's password to default"""
    user = User.query.get_or_404(id)
    user.set_password('password123')
    db.session.commit()
    flash(f'Password for {user.get_display_name()} has been reset to "password123".', 'success')
    return redirect(url_for('users'))


@app.route('/users/<int:id>/delete', methods=['POST'])
@admin_required
def delete_user(id):
    """Delete a user"""
    user = User.query.get_or_404(id)
    
    # Prevent deleting the last admin or team leader
    admin_and_team_leader_count = User.query.filter(User.role.in_(['admin', 'team_leader'])).count()
    if user.is_admin() and admin_and_team_leader_count == 1:
        flash('Cannot delete the last admin or team leader user.', 'danger')
        return redirect(url_for('users'))
    
    # Prevent deleting yourself
    if user.id == current_user.id:
        flash('You cannot delete your own account.', 'danger')
        return redirect(url_for('users'))
    
    # Get an admin or team leader user to reassign created_by fields (use current_user if admin/team leader, otherwise find any admin/team leader)
    admin_user = current_user if current_user.is_admin() else User.query.filter(User.role.in_(['admin', 'team_leader'])).first()
    if not admin_user:
        admin_user = current_user  # Fallback to current user
    
    # Delete all records where user is directly referenced (NOT NULL constraints)
    # 1. UserPermission records
    UserPermission.query.filter_by(user_id=user.id).delete()
    UserPermission.query.filter_by(granted_by=user.id).delete()
    
    # 2. Post interactions
    PostLike.query.filter_by(user_id=user.id).delete()
    PostHeart.query.filter_by(user_id=user.id).delete()
    PostRepost.query.filter_by(user_id=user.id).delete()
    PostComment.query.filter_by(user_id=user.id).delete()
    
    # 3. Notifications
    Notification.query.filter_by(user_id=user.id).delete()
    Notification.query.filter_by(actor_id=user.id).delete()
    
    # 4. Journal entries
    Journal.query.filter_by(user_id=user.id).delete()
    
    # Reassign created_by fields to admin (to preserve data)
    # 5. SundayService
    SundayService.query.filter_by(created_by=user.id).update({'created_by': admin_user.id})
    
    # 6. Practice
    Practice.query.filter_by(created_by=user.id).update({'created_by': admin_user.id})
    
    # 7. Song
    Song.query.filter_by(created_by=user.id).update({'created_by': admin_user.id})
    
    # 8. Slide
    Slide.query.filter_by(created_by=user.id).update({'created_by': admin_user.id})
    
    # 9. EventAnnouncement
    EventAnnouncement.query.filter_by(created_by=user.id).update({'created_by': admin_user.id})
    
    # Note: PracticeSong.prepared_by is nullable, so no action needed
    # Note: SMSLog fields are nullable, so no action needed
    # Note: Musician.user_id is nullable, so no action needed
    
    # Now delete the user
    db.session.delete(user)
    db.session.commit()
    flash('User deleted successfully. All related records have been handled.', 'success')
    return redirect(url_for('users'))


# Initialize database
def migrate_database():
    """Add missing columns to existing database"""
    try:
        from sqlalchemy import text  # type: ignore
        with db.engine.connect() as conn:
            # Migrate song table - add gender_key column
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='song'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(song)"))
                columns = [row[1] for row in result]
                
                if 'gender_key' not in columns:
                    print('Adding gender_key column to song table...')
                    conn.execute(text("ALTER TABLE song ADD COLUMN gender_key VARCHAR(10)"))
                    conn.commit()
                    print('Migration completed: gender_key column added')
            
            # Create leave_request table if it doesn't exist
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='leave_request'"
            ))
            table_exists = result.fetchone() is not None
            
            if not table_exists:
                print('Creating leave_request table...')
                conn.execute(text("""
                    CREATE TABLE leave_request (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        user_id INTEGER NOT NULL,
                        musician_id INTEGER NOT NULL,
                        date DATE NOT NULL,
                        reason VARCHAR(500) NOT NULL,
                        status VARCHAR(20) NOT NULL DEFAULT 'pending',
                        requested_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        reviewed_by INTEGER,
                        reviewed_at DATETIME,
                        review_notes VARCHAR(500),
                        FOREIGN KEY (user_id) REFERENCES user (id),
                        FOREIGN KEY (musician_id) REFERENCES musician (id),
                        FOREIGN KEY (reviewed_by) REFERENCES user (id)
                    )
                """))
                conn.commit()
                print('Migration completed: leave_request table created')
            
            # Migrate practice table - add purpose column
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='practice'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(practice)"))
                columns = [row[1] for row in result]
                
                if 'purpose' not in columns:
                    print('Adding purpose column to practice table...')
                    conn.execute(text("ALTER TABLE practice ADD COLUMN purpose VARCHAR(200)"))
                    conn.commit()
                    print('Migration completed: purpose column added')
            
            # Migrate slide table - add language column if needed
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='slide'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(slide)"))
                columns = [row[1] for row in result]
                
                # Add description column if it doesn't exist
                if 'description' not in columns:
                    print('Adding description column to slide table...')
                    conn.execute(text("ALTER TABLE slide ADD COLUMN description TEXT"))
                    conn.commit()
                    print('Migration completed: description column added')
            
            # Migrate notification table - add leave_request_id column
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='notification'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(notification)"))
                columns = [row[1] for row in result]
                
                if 'leave_request_id' not in columns:
                    print('Adding leave_request_id column to notification table...')
                    conn.execute(text("ALTER TABLE notification ADD COLUMN leave_request_id INTEGER"))
                    conn.commit()
                    print('Migration completed: leave_request_id column added')
                
                # Add language column if it doesn't exist
                if 'language' not in columns:
                    print('Adding language column to slide table...')
                    conn.execute(text("ALTER TABLE slide ADD COLUMN language VARCHAR(20)"))
                    conn.commit()
                    print('Migration completed: language column added')
                
                # Add file_type column if it doesn't exist
                if 'file_type' not in columns:
                    print('Adding file_type column to slide table...')
                    conn.execute(text("ALTER TABLE slide ADD COLUMN file_type VARCHAR(50)"))
                    conn.commit()
                    print('Migration completed: file_type column added')
            
            # Migrate user table - add nickname and mobile_number columns
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='user'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(user)"))
                columns = [row[1] for row in result]
                
                if 'nickname' not in columns:
                    print('Adding nickname column to user table...')
                    conn.execute(text("ALTER TABLE user ADD COLUMN nickname VARCHAR(100)"))
                    conn.commit()
                    print('Migration completed: nickname column added')
                
                if 'mobile_number' not in columns:
                    print('Adding mobile_number column to user table...')
                    conn.execute(text("ALTER TABLE user ADD COLUMN mobile_number VARCHAR(20)"))
                    conn.commit()
                    print('Migration completed: mobile_number column added')
            
            # Migrate musician table - add bio, roles, interests columns
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='musician'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(musician)"))
                columns = [row[1] for row in result]
                
                if 'bio' not in columns:
                    print('Adding bio column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN bio TEXT"))
                    conn.commit()
                    print('Migration completed: bio column added')
                
                if 'roles' not in columns:
                    print('Adding roles column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN roles VARCHAR(200)"))
                    conn.commit()
                    print('Migration completed: roles column added')
                
                if 'interests' not in columns:
                    print('Adding interests column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN interests VARCHAR(300)"))
                    conn.commit()
                    print('Migration completed: interests column added')
                
                if 'profile_picture' not in columns:
                    print('Adding profile_picture column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN profile_picture VARCHAR(255)"))
                    conn.commit()
                    print('Migration completed: profile_picture column added')
                
                if 'banner' not in columns:
                    print('Adding banner column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN banner VARCHAR(255)"))
                    conn.commit()
                    print('Migration completed: banner column added')
                
                if 'mobile' not in columns:
                    print('Adding mobile column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN mobile VARCHAR(20)"))
                    conn.commit()
                    print('Migration completed: mobile column added')
                
                if 'outlook_email' not in columns:
                    print('Adding outlook_email column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN outlook_email VARCHAR(120)"))
                    conn.commit()
                    print('Migration completed: outlook_email column added')
                
                if 'whatsapp' not in columns:
                    print('Adding whatsapp column to musician table...')
                    conn.execute(text("ALTER TABLE musician ADD COLUMN whatsapp VARCHAR(20)"))
                    conn.commit()
                    print('Migration completed: whatsapp column added')
            
            # Create profile_post table if it doesn't exist
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='profile_post'"
            ))
            table_exists = result.fetchone() is not None
            
            if not table_exists:
                print('Creating profile_post table...')
                conn.execute(text("""
                    CREATE TABLE profile_post (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        musician_id INTEGER NOT NULL,
                        content TEXT,
                        image_path VARCHAR(255),
                        video_path VARCHAR(255),
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (musician_id) REFERENCES musician (id)
                    )
                """))
                conn.commit()
                print('Migration completed: profile_post table created')
            
            # Create post interaction tables if they don't exist
            interaction_tables = [
                ('post_like', """
                    CREATE TABLE post_like (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        post_id INTEGER NOT NULL,
                        user_id INTEGER NOT NULL,
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (post_id) REFERENCES profile_post (id),
                        FOREIGN KEY (user_id) REFERENCES user (id),
                        UNIQUE(post_id, user_id)
                    )
                """),
                ('post_heart', """
                    CREATE TABLE post_heart (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        post_id INTEGER NOT NULL,
                        user_id INTEGER NOT NULL,
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (post_id) REFERENCES profile_post (id),
                        FOREIGN KEY (user_id) REFERENCES user (id),
                        UNIQUE(post_id, user_id)
                    )
                """),
                ('post_repost', """
                    CREATE TABLE post_repost (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        post_id INTEGER NOT NULL,
                        user_id INTEGER NOT NULL,
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (post_id) REFERENCES profile_post (id),
                        FOREIGN KEY (user_id) REFERENCES user (id),
                        UNIQUE(post_id, user_id)
                    )
                """),
                ('post_comment', """
                    CREATE TABLE post_comment (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        post_id INTEGER NOT NULL,
                        user_id INTEGER NOT NULL,
                        content TEXT NOT NULL,
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (post_id) REFERENCES profile_post (id),
                        FOREIGN KEY (user_id) REFERENCES user (id)
                    )
                """)
            ]
            
            for table_name, create_sql in interaction_tables:
                result = conn.execute(text(
                    f"SELECT name FROM sqlite_master WHERE type='table' AND name='{table_name}'"
                ))
                table_exists = result.fetchone() is not None
                
                if not table_exists:
                    print(f'Creating {table_name} table...')
                    conn.execute(text(create_sql))
                    conn.commit()
                    print(f'Migration completed: {table_name} table created')
            
            # Create practice_song table if it doesn't exist
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='practice_song'"
            ))
            table_exists = result.fetchone() is not None
            
            if not table_exists:
                print('Creating practice_song table...')
                conn.execute(text("""
                    CREATE TABLE practice_song (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        practice_id INTEGER NOT NULL,
                        song_id INTEGER,
                        song_name VARCHAR(200),
                        speed VARCHAR(20),
                        prepared_by INTEGER,
                        "order" INTEGER DEFAULT 0,
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (practice_id) REFERENCES practice (id),
                        FOREIGN KEY (song_id) REFERENCES song (id),
                        FOREIGN KEY (prepared_by) REFERENCES user (id)
                    )
                """))
                conn.commit()
                print('Migration completed: practice_song table created')
            else:
                # Check if song_id is nullable (SQLite doesn't support ALTER COLUMN, so we need to recreate)
                result = conn.execute(text("PRAGMA table_info(practice_song)"))
                columns = {row[1]: row for row in result.fetchall()}
                
                # Check if song_id has NOT NULL constraint (column[3] is notnull flag)
                song_id_info = columns.get('song_id')
                needs_recreate = False
                
                if song_id_info and song_id_info[3] == 1:  # NOT NULL constraint exists
                    print('song_id has NOT NULL constraint, recreating table to make it nullable...')
                    needs_recreate = True
                
                # Check for missing columns
                if 'song_name' not in columns or 'speed' not in columns or 'prepared_by' not in columns:
                    needs_recreate = True
                
                if needs_recreate:
                    # Backup existing data
                    print('Backing up existing practice_song data...')
                    backup_data = conn.execute(text("SELECT * FROM practice_song")).fetchall()
                    
                    # Drop old table
                    conn.execute(text("DROP TABLE practice_song"))
                    conn.commit()
                    
                    # Create new table with correct schema
                    print('Creating new practice_song table with nullable song_id...')
                    conn.execute(text("""
                        CREATE TABLE practice_song (
                            id INTEGER PRIMARY KEY AUTOINCREMENT,
                            practice_id INTEGER NOT NULL,
                            song_id INTEGER,
                            song_name VARCHAR(200),
                            "key" VARCHAR(20),
                            speed VARCHAR(20),
                            prepared_by INTEGER,
                            "order" INTEGER DEFAULT 0,
                            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                            FOREIGN KEY (practice_id) REFERENCES practice (id),
                            FOREIGN KEY (song_id) REFERENCES song (id),
                            FOREIGN KEY (prepared_by) REFERENCES user (id)
                        )
                    """))
                    conn.commit()
                    
                    # Restore data (only if there was data)
                    if backup_data:
                        print('Restoring practice_song data...')
                        for row in backup_data:
                            # Map old columns to new structure
                            # Assuming old structure: id, practice_id, song_id, order, created_at
                            old_id = row[0]
                            old_practice_id = row[1]
                            old_song_id = row[2] if len(row) > 2 else None
                            old_order = row[3] if len(row) > 3 else 0
                            old_created_at = row[4] if len(row) > 4 else None
                            
                            conn.execute(text("""
                                INSERT INTO practice_song (id, practice_id, song_id, "order", created_at)
                                VALUES (?, ?, ?, ?, ?)
                            """), (old_id, old_practice_id, old_song_id, old_order, old_created_at))
                        conn.commit()
                    
                    print('Migration completed: practice_song table recreated with nullable song_id')
                else:
                    # Just add missing columns if they don't exist
                    if 'song_name' not in columns:
                        print('Adding song_name column to practice_song table...')
                        conn.execute(text("ALTER TABLE practice_song ADD COLUMN song_name VARCHAR(200)"))
                        conn.commit()
                        print('Migration completed: song_name column added')
                    
                    if 'key' not in columns:
                        print('Adding key column to practice_song table...')
                        conn.execute(text('ALTER TABLE practice_song ADD COLUMN "key" VARCHAR(20)'))
                        conn.commit()
                        print('Migration completed: key column added')
                    
                    if 'speed' not in columns:
                        print('Adding speed column to practice_song table...')
                        conn.execute(text("ALTER TABLE practice_song ADD COLUMN speed VARCHAR(20)"))
                        conn.commit()
                        print('Migration completed: speed column added')
                    
                    if 'prepared_by' not in columns:
                        print('Adding prepared_by column to practice_song table...')
                        conn.execute(text("ALTER TABLE practice_song ADD COLUMN prepared_by INTEGER"))
                        conn.execute(text("""
                            CREATE INDEX IF NOT EXISTS idx_practice_song_prepared_by 
                            ON practice_song(prepared_by)
                        """))
                        conn.commit()
                        print('Migration completed: prepared_by column added')
            
            # Migrate tool table - add developer_name column
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='tool'"
            ))
            table_exists = result.fetchone() is not None
            
            if table_exists:
                result = conn.execute(text("PRAGMA table_info(tool)"))
                columns = [row[1] for row in result.fetchall()]
                
                if 'developer_name' not in columns:
                    print('Adding developer_name column to tool table...')
                    with conn.begin():
                        conn.execute(text("ALTER TABLE tool ADD COLUMN developer_name VARCHAR(200)"))
                    print('Migration completed: developer_name column added')
                else:
                    print('developer_name column already exists in tool table')
            
            # Create message table for chat if it doesn't exist
            result = conn.execute(text(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='message'"
            ))
            table_exists = result.fetchone() is not None
            
            if not table_exists:
                print('Creating message table for chat...')
                conn.execute(text("""
                    CREATE TABLE message (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        user_id INTEGER NOT NULL,
                        recipient_id INTEGER,
                        content TEXT NOT NULL,
                        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                        FOREIGN KEY (user_id) REFERENCES user (id),
                        FOREIGN KEY (recipient_id) REFERENCES user (id)
                    )
                """))
                conn.commit()
                print('Migration completed: message table created')
            else:
                # Check if recipient_id column exists, add it if not
                result = conn.execute(text("PRAGMA table_info(message)"))
                columns = [row[1] for row in result.fetchall()]
                
                if 'recipient_id' not in columns:
                    print('Adding recipient_id column to message table...')
                    conn.execute(text("ALTER TABLE message ADD COLUMN recipient_id INTEGER"))
                    conn.commit()
                    print('Migration completed: recipient_id column added')
                else:
                    print('message table already has recipient_id column')
            
            # Handle password_hash NULL values for existing databases
                # SQLite doesn't support changing column nullability, so we ensure all users have a password_hash
                try:
                    from werkzeug.security import generate_password_hash  # type: ignore
                    default_hash = generate_password_hash('password123')
                    result = conn.execute(text("SELECT COUNT(*) FROM user WHERE password_hash IS NULL"))
                    null_count = result.fetchone()[0]
                    if null_count > 0:
                        conn.execute(text("UPDATE user SET password_hash = :hash WHERE password_hash IS NULL"),
                                   {'hash': default_hash})
                        conn.commit()
                        print(f'Migration: Updated {null_count} NULL password_hash values to default')
                except Exception as e:
                    print(f'Migration note for password_hash: {e}')
    except Exception as e:
        # Table might not exist yet, which is fine - it will be created by create_all()
        print(f'Migration check: {e}')


def init_db():
    """Create database tables and initial admin user if needed"""
    with app.app_context():
        db.create_all()
        
        # Run migrations for existing databases
        migrate_database()
        
        # Create default admin user if no users exist
        if User.query.count() == 0:
            admin = User(
                username='admin',
                email='admin@example.com',
                role='admin'
            )
            admin.set_password('admin123')  # Set default password
            db.session.add(admin)
            db.session.commit()
            print('Default admin user created: username=admin, password=admin123')


# Notification routes
@app.route('/notifications')
@login_required
def get_notifications():
    """Get notifications for current user"""
    notifications = Notification.query.filter_by(user_id=current_user.id).order_by(Notification.created_at.desc()).limit(20).all()
    
    notifications_data = []
    for notif in notifications:
        actor_name = notif.actor.get_display_name() if notif.actor else 'Someone'
        notification_text = ''
        link = '#'
        
        if notif.notification_type == 'like':
            notification_text = f"{actor_name} liked your post"
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'heart':
            notification_text = f"{actor_name} ❤️ your post"
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'share':
            notification_text = f"{actor_name} shared your post"
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'comment':
            notification_text = f"{actor_name} commented on your post"
            if notif.post:
                link = url_for('view_musician_profile', id=notif.post.musician_id)
        elif notif.notification_type == 'practice':
            notification_text = f"{actor_name} created a new practice schedule"
            if notif.practice:
                link = url_for('practice_detail', id=notif.practice.id)
        elif notif.notification_type == 'leave_request':
            notification_text = f"{actor_name} filed a leave request for your approval"
            link = '#'  # Will be handled by JavaScript to show popup
        elif notif.notification_type == 'leave_approved':
            notification_text = f"Your leave request has been approved by {actor_name}"
            link = url_for('leave_requests')
        elif notif.notification_type == 'leave_rejected':
            # Include rejection reason if available
            if notif.leave_request and notif.leave_request.review_notes:
                notification_text = f"Your leave request has been rejected by {actor_name}. Reason: {notif.leave_request.review_notes}"
            else:
                notification_text = f"Your leave request has been rejected by {actor_name}"
            link = url_for('leave_requests')
        
        notification_data = {
            'id': notif.id,
            'text': notification_text,
            'link': link,
            'is_read': notif.is_read,
            'created_at': notif.created_at.strftime('%Y-%m-%d %H:%M:%S'),
            'time_ago': _time_ago(notif.created_at),
            'type': notif.notification_type
        }
        
        # Add leave request ID if applicable
        if notif.leave_request_id:
            notification_data['leave_request_id'] = notif.leave_request_id
        
        notifications_data.append(notification_data)
    
    unread_count = Notification.query.filter_by(user_id=current_user.id, is_read=False).count()
    
    return jsonify({
        'notifications': notifications_data,
        'unread_count': unread_count
    })


def _time_ago(dt):
    """Calculate time ago string"""
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc)
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    diff = now - dt
    
    if diff.days > 0:
        return f"{diff.days} day{'s' if diff.days > 1 else ''} ago"
    elif diff.seconds >= 3600:
        hours = diff.seconds // 3600
        return f"{hours} hour{'s' if hours > 1 else ''} ago"
    elif diff.seconds >= 60:
        minutes = diff.seconds // 60
        return f"{minutes} minute{'s' if minutes > 1 else ''} ago"
    else:
        return "Just now"


@app.route('/notifications/<int:notification_id>/read', methods=['POST'])
@login_required
def mark_notification_read(notification_id):
    """Mark a notification as read"""
    notification = Notification.query.get_or_404(notification_id)
    
    # Ensure user owns this notification
    if notification.user_id != current_user.id:
        return jsonify({'success': False, 'message': 'Unauthorized'}), 403
    
    # Mark as read
    notification.is_read = True
    try:
        db.session.commit()
        # Refresh to ensure the change is persisted
        db.session.refresh(notification)
        return jsonify({'success': True, 'is_read': notification.is_read})
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': str(e)}), 500


@app.route('/notifications/read-all', methods=['POST'])
@login_required
def mark_all_notifications_read():
    """Mark all notifications as read for current user"""
    Notification.query.filter_by(user_id=current_user.id, is_read=False).update({'is_read': True})
    db.session.commit()
    
    flash('All notifications marked as read.', 'success')
    return redirect(url_for('notifications_page'))


@app.route('/notifications/delete-all', methods=['POST'])
@csrf.exempt
@login_required
def delete_all_notifications():
    """Delete all notifications for current user"""
    try:
        deleted_count = Notification.query.filter_by(user_id=current_user.id).delete()
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': f'Successfully deleted {deleted_count} notification(s).',
            'count': deleted_count
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'message': f'Error deleting notifications: {str(e)}'}), 500


def schedule_practice_sms_reminders(practice, musician):
    """
    Schedule SMS reminders for a practice (1 day before and 1 hour before)
    
    Args:
        practice: Practice object
        musician: Musician object
    """
    if not practice.date or not practice.time:
        return
    
    # Combine date and time
    practice_datetime = datetime.combine(practice.date, practice.time)
    
    # Calculate reminder times
    one_day_before = practice_datetime - timedelta(days=1)
    one_hour_before = practice_datetime - timedelta(hours=1)
    
    # Only schedule if reminders are in the future
    now = datetime.now()
    
    # Check if scheduler is available
    if scheduler is None or not scheduler.running:
        return
    
    if one_day_before > now:
        try:
            scheduler.add_job(
                func=send_reminder_sms_job,
                trigger=DateTrigger(run_date=one_day_before),
                args=[practice.id, musician.id, 'day_before'],
                id=f'practice_{practice.id}_musician_{musician.id}_day_before',
                replace_existing=True
            )
        except Exception as e:
            print(f"Warning: Could not schedule day_before reminder: {e}")
    
    if one_hour_before > now:
        try:
            scheduler.add_job(
                func=send_reminder_sms_job,
                trigger=DateTrigger(run_date=one_hour_before),
                args=[practice.id, musician.id, 'hour_before'],
                id=f'practice_{practice.id}_musician_{musician.id}_hour_before',
                replace_existing=True
            )
        except Exception as e:
            print(f"Warning: Could not schedule hour_before reminder: {e}")


def send_reminder_sms_job(practice_id, musician_id, reminder_type):
    """
    Background job to send SMS reminder
    
    Args:
        practice_id: Practice ID
        musician_id: Musician ID
        reminder_type: 'day_before' or 'hour_before'
    """
    with app.app_context():
        practice = Practice.query.get(practice_id)
        musician = Musician.query.get(musician_id)
        
        if practice and musician:
            result = send_practice_reminder_sms(practice, musician, reminder_type)
            # Handle old (success, error), new (success, error, sid), and latest (success, error, sid, status) formats
            if len(result) == 4:
                success, error, message_sid, twilio_status = result
            elif len(result) == 3:
                success, error, message_sid = result
                twilio_status = None
            else:
                success, error = result[0], result[1] if len(result) > 1 else None
                message_sid = None
                twilio_status = None
            
            # Get user info for logging
            user = musician.user if musician.user_id else None
            recipient_name = user.get_display_name() if user else musician.get_display_name()
            recipient_phone = format_phone_number(user.mobile_number) if user and user.mobile_number else None
            
            # Log SMS attempt
            try:
                sms_log = SMSLog(
                    recipient_user_id=user.id if user else None,
                    recipient_phone=recipient_phone or 'Unknown',
                    recipient_name=recipient_name,
                    message_type=f'practice_reminder_{reminder_type}',
                    practice_id=practice.id,
                    musician_id=musician.id,
                    message_content=f"Practice reminder ({reminder_type}) for {practice.date.strftime('%B %d, %Y') if practice.date else 'TBD'}",
                    status='success' if success else 'failed',
                    twilio_status=twilio_status,
                    error_message=error if not success else None,
                    twilio_message_sid=message_sid,
                    sent_by_user_id=None  # System-scheduled, no user trigger
                )
                db.session.add(sms_log)
                db.session.commit()
            except Exception as log_error:
                print(f"Warning: Could not log SMS reminder: {log_error}")


@app.route('/api/bible-verse', methods=['POST'])
@login_required
def fetch_bible_verse():
    """Fetch Bible verse content from an API"""
    try:
        data = request.get_json()
        book = data.get('book', '')
        chapter = data.get('chapter', 0)
        verse = data.get('verse', 0)
        version = data.get('version', 'NIV')
        
        if not book or not chapter or not verse:
            return jsonify({'error': 'Missing book, chapter, or verse'}), 400
        
        # Use Bible Gateway API or alternative
        # Format: book name, chapter:verse
        import urllib.parse
        book_encoded = urllib.parse.quote(book)
        passage = f"{book_encoded}+{chapter}:{verse}"
        
        # Try to fetch from Bible Gateway (note: may require API key for production)
        # For now, we'll use a simple web scraping approach or return formatted reference
        # In production, you'd want to use a proper Bible API service
        
        # Using Bible Gateway's public URL (may have CORS issues, so we'll parse server-side)
        try:
            import requests  # type: ignore
            from bs4 import BeautifulSoup  # type: ignore
            
            url = f"https://www.biblegateway.com/passage/?search={passage}&version={version}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers, timeout=5)
            
            if response.status_code == 200:
                soup = BeautifulSoup(response.content, 'html.parser')
                # Find the verse content
                verse_div = soup.find('div', class_='verse')
                if verse_div:
                    # Get text content, removing verse numbers
                    content = verse_div.get_text(strip=True)
                    # Remove verse number if present
                    content = content.replace(f'{verse}', '').strip()
                    return jsonify({'content': content})
        except Exception as e:
            print(f"Error fetching from Bible Gateway: {e}")
        
        # Fallback: return formatted reference
        return jsonify({
            'content': f"{book} {chapter}:{verse} - [Content not available. Please check the reference.]"
        })
        
    except Exception as e:
        print(f"Error in fetch_bible_verse: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/chat', methods=['GET'])
@login_required
def chat():
    """Chat page for team communication"""
    # Get recent messages (last 100)
    messages = Message.query.order_by(Message.created_at.desc()).limit(100).all()
    messages.reverse()  # Show oldest first
    
    # Get all users for display
    users = User.query.all()
    user_dict = {user.id: user for user in users}
    
    return render_template('chat.html', messages=messages, user_dict=user_dict)


@app.route('/api/chat/messages', methods=['GET'])
@csrf.exempt
@login_required
def get_messages():
    """API endpoint to get messages (for polling) - group chat only"""
    try:
        recipient_id = request.args.get('recipient_id')
        
        # If recipient_id is provided, get private messages
        if recipient_id:
            try:
                recipient_id = int(recipient_id)
                # Get messages between current user and recipient
                messages = Message.query.filter(
                    ((Message.user_id == current_user.id) & (Message.recipient_id == recipient_id)) |
                    ((Message.user_id == recipient_id) & (Message.recipient_id == current_user.id))
                ).order_by(Message.created_at.asc()).all()
            except:
                messages = []
        else:
            # Get group chat messages (recipient_id is None)
            since = request.args.get('since')
            if since:
                try:
                    since_dt = datetime.fromisoformat(since.replace('Z', '+00:00'))
                    messages = Message.query.filter(
                        Message.recipient_id.is_(None),
                        Message.created_at > since_dt
                    ).order_by(Message.created_at.asc()).all()
                except:
                    messages = Message.query.filter(Message.recipient_id.is_(None)).order_by(Message.created_at.desc()).limit(50).all()
                    messages.reverse()
            else:
                # Get last 50 group messages
                messages = Message.query.filter(Message.recipient_id.is_(None)).order_by(Message.created_at.desc()).limit(50).all()
                messages.reverse()
        
        # Format messages for JSON
        messages_data = []
        for msg in messages:
            user = User.query.get(msg.user_id)
            profile_picture = None
            if user and user.musician and user.musician.profile_picture:
                profile_picture = user.musician.profile_picture
            
            messages_data.append({
                'id': msg.id,
                'user_id': msg.user_id,
                'recipient_id': msg.recipient_id,
                'username': user.get_display_name() if user else 'Unknown',
                'profile_picture': profile_picture,
                'content': msg.content,
                'created_at': msg.created_at.isoformat(),
                'is_own': msg.user_id == current_user.id
            })
        
        return jsonify({
            'success': True,
            'messages': messages_data,
            'latest_timestamp': messages[-1].created_at.isoformat() if messages else None
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/chat/send', methods=['POST'])
@csrf.exempt
@login_required
def send_message():
    """API endpoint to send a message (group or private)"""
    try:
        data = request.get_json()
        content = data.get('content', '').strip()
        recipient_id = data.get('recipient_id')  # None for group chat, user_id for private
        
        if not content:
            return jsonify({'success': False, 'error': 'Message cannot be empty'}), 400
        
        # Validate recipient if provided
        if recipient_id:
            try:
                recipient_id = int(recipient_id)
                recipient = User.query.get(recipient_id)
                if not recipient:
                    return jsonify({'success': False, 'error': 'Recipient not found'}), 404
                if recipient_id == current_user.id:
                    return jsonify({'success': False, 'error': 'Cannot send message to yourself'}), 400
            except:
                return jsonify({'success': False, 'error': 'Invalid recipient'}), 400
        
        # Create message
        message = Message(
            user_id=current_user.id,
            recipient_id=recipient_id if recipient_id else None,
            content=content
        )
        db.session.add(message)
        db.session.commit()
        
        # Return the created message
        user = User.query.get(current_user.id)
        profile_picture = None
        if user and user.musician and user.musician.profile_picture:
            profile_picture = user.musician.profile_picture
        
        return jsonify({
            'success': True,
            'message': {
                'id': message.id,
                'user_id': message.user_id,
                'recipient_id': message.recipient_id,
                'username': user.get_display_name() if user else 'Unknown',
                'profile_picture': profile_picture,
                'content': message.content,
                'created_at': message.created_at.isoformat(),
                'is_own': True
            }
        })
    except Exception as e:
        # If error is due to missing recipient_id column, run migration and retry
        if 'recipient_id' in str(e) or 'no such column' in str(e).lower():
            try:
                # Run migration to add recipient_id column
                from sqlalchemy import text
                with db.engine.connect() as conn:
                    result = conn.execute(text("PRAGMA table_info(message)"))
                    columns = [row[1] for row in result.fetchall()]
                    
                    if 'recipient_id' not in columns:
                        print('Adding recipient_id column to message table...')
                        conn.execute(text("ALTER TABLE message ADD COLUMN recipient_id INTEGER"))
                        conn.commit()
                        print('Migration completed: recipient_id column added')
                
                # Retry creating the message
                message = Message(
                    user_id=current_user.id,
                    recipient_id=recipient_id if recipient_id else None,
                    content=content
                )
                db.session.add(message)
                db.session.commit()
                
                # Return the created message
                user = User.query.get(current_user.id)
                profile_picture = None
                if user and user.musician and user.musician.profile_picture:
                    profile_picture = user.musician.profile_picture
                
                return jsonify({
                    'success': True,
                    'message': {
                        'id': message.id,
                        'user_id': message.user_id,
                        'recipient_id': message.recipient_id,
                        'username': user.get_display_name() if user else 'Unknown',
                        'profile_picture': profile_picture,
                        'content': message.content,
                        'created_at': message.created_at.isoformat(),
                        'is_own': True
                    }
                })
            except Exception as migration_error:
                return jsonify({'success': False, 'error': f'Database error: {str(migration_error)}'}), 500
        else:
            return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/chat/users', methods=['GET'])
@csrf.exempt
@login_required
def get_chat_users():
    """API endpoint to get all users for the members list"""
    try:
        users = User.query.order_by(User.username).all()
        users_data = []
        for user in users:
            # Get unread count for private messages from this user
            unread_count = Message.query.filter(
                Message.user_id == user.id,
                Message.recipient_id == current_user.id
            ).count()
            
            users_data.append({
                'id': user.id,
                'username': user.get_display_name(),
                'role': user.role,
                'unread_count': unread_count
            })
        
        return jsonify({
            'success': True,
            'users': users_data
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/chat/conversations', methods=['GET'])
@csrf.exempt
@login_required
def get_conversations():
    """API endpoint to get list of conversations"""
    try:
        # Get all users the current user has messaged or been messaged by
        sent_messages = Message.query.filter_by(user_id=current_user.id).filter(Message.recipient_id.isnot(None)).all()
        received_messages = Message.query.filter_by(recipient_id=current_user.id).all()
        
        # Get unique user IDs
        user_ids = set()
        for msg in sent_messages:
            if msg.recipient_id:
                user_ids.add(msg.recipient_id)
        for msg in received_messages:
            user_ids.add(msg.user_id)
        
        # Get conversation details
        conversations = []
        for user_id in user_ids:
            user = User.query.get(user_id)
            if not user:
                continue
            
            # Get last message in this conversation
            last_message = Message.query.filter(
                ((Message.user_id == current_user.id) & (Message.recipient_id == user_id)) |
                ((Message.user_id == user_id) & (Message.recipient_id == current_user.id))
            ).order_by(Message.created_at.desc()).first()
            
            # Count unread messages (messages from other user after last viewed)
            unread_count = Message.query.filter(
                Message.user_id == user_id,
                Message.recipient_id == current_user.id,
                Message.created_at > (last_message.created_at if last_message else datetime.min)
            ).count()
            
            conversations.append({
                'user_id': user_id,
                'username': user.get_display_name(),
                'last_message': last_message.content if last_message else '',
                'last_message_time': last_message.created_at.isoformat() if last_message else None,
                'unread_count': unread_count
            })
        
        # Sort by last message time
        conversations.sort(key=lambda x: x['last_message_time'] or '', reverse=True)
        
        return jsonify({
            'success': True,
            'conversations': conversations
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/migrate', methods=['GET'])
@admin_required
def run_migration():
    """Manually trigger database migration"""
    try:
        migrate_database()
        flash('Migration completed successfully!', 'success')
    except Exception as e:
        flash(f'Migration error: {str(e)}', 'danger')
    return redirect(url_for('dashboard'))

if __name__ == '__main__':
    # Ensure directories exist
    os.makedirs(os.path.join(app.root_path, 'instance'), exist_ok=True)
    os.makedirs(os.path.join(app.root_path, 'static', 'chords'), exist_ok=True)
    os.makedirs(os.path.join(app.root_path, 'static', 'slides'), exist_ok=True)
    os.makedirs(os.path.join(app.root_path, 'static', 'announcements'), exist_ok=True)
    
    # Initialize database
    init_db()
    
    app.run(host='0.0.0.0', port=5000, debug=True)

